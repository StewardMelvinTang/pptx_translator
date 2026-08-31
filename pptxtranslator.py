import os
import ssl
import subprocess
import sys
import hashlib
from datetime import datetime, timezone
from pathlib import Path

# Fix conda SSL_CERT_FILE pointing to nonexistent path
if "SSL_CERT_FILE" in os.environ and not os.path.exists(os.environ["SSL_CERT_FILE"]):
    del os.environ["SSL_CERT_FILE"]
if "SSL_CERT_DIR" in os.environ and not os.path.exists(os.environ["SSL_CERT_DIR"]):
    del os.environ["SSL_CERT_DIR"]

import customtkinter as ctk
import tkinter as tk
from tkinter import filedialog, messagebox
from PIL import Image
try:
    from tkinterdnd2 import DND_FILES, TkinterDnD
except ImportError:
    DND_FILES = None
    TkinterDnD = None
from pptx import Presentation
import json
import time
import threading
import re
import base64
import html
import webbrowser
import pymupdf as fitz
from openai import OpenAI

APP_NAME = "Document Translator"
APP_DISPLAY_NAME = "PPT Translator"
APP_VERSION = "4.0.0"

IS_WINDOWS = sys.platform.startswith("win")
IS_MACOS = sys.platform == "darwin"

if IS_MACOS:
    UI_FONT_FAMILY = "SF Pro Text"
elif IS_WINDOWS:
    UI_FONT_FAMILY = "Segoe UI Variable"
else:
    UI_FONT_FAMILY = "DejaVu Sans"

# A deliberately neutral palette that reads well beside both macOS and Windows chrome.
COLORS = {
    "background": ("#F4F6FA", "#0E1014"),
    "surface": ("#FFFFFF", "#181B21"),
    "surface_alt": ("#F7F8FB", "#20242B"),
    "border": ("#E1E5EC", "#303640"),
    "text": ("#171A21", "#F3F5F8"),
    "muted": ("#687080", "#9CA5B4"),
    "accent": ("#5367F8", "#7180FF"),
    "accent_hover": ("#4053E4", "#8290FF"),
    "danger": ("#D9485F", "#FF7185"),
    "success": ("#1A9B68", "#48C78E"),
}

_ICON_CACHE = {}


def ui_font(size=13, weight="normal"):
    return ctk.CTkFont(family=UI_FONT_FAMILY, size=size, weight=weight)


def resolve_color(color):
    """Resolve a CustomTkinter light/dark tuple for native Tk widgets."""
    if isinstance(color, (tuple, list)):
        return color[1] if ctk.get_appearance_mode() == "Dark" else color[0]
    return color


def load_icon(name, size=18):
    """Load a theme-aware Lucide icon from bundled raster assets."""
    key = (name, size)
    if key not in _ICON_CACHE:
        light_path = resource_path(f"assets/icons/{name}-light.png")
        dark_path = resource_path(f"assets/icons/{name}-dark.png")
        _ICON_CACHE[key] = ctk.CTkImage(
            light_image=Image.open(light_path),
            dark_image=Image.open(dark_path),
            size=(size, size),
        )
    return _ICON_CACHE[key]


def load_brand_image(size=30):
    key = ("brand", size)
    if key not in _ICON_CACHE:
        image = Image.open(resource_path("pptx_icon.png"))
        _ICON_CACHE[key] = ctk.CTkImage(
            light_image=image, dark_image=image, size=(size, size)
        )
    return _ICON_CACHE[key]


def ellipsize_middle(text, max_chars):
    """Keep long file paths from forcing the window wider than the display."""
    if len(text) <= max_chars:
        return text
    side = max(1, (max_chars - 1) // 2)
    return f"{text[:side]}\u2026{text[-side:]}"


def resource_path(filename):
    """Return an asset path that works in source and PyInstaller bundles."""
    bundle_dir = Path(getattr(sys, "_MEIPASS", Path(__file__).resolve().parent))
    return bundle_dir / filename


def get_config_path():
    """Use the native per-user application data directory on each OS."""
    if IS_WINDOWS:
        base = Path(os.environ.get("APPDATA", Path.home() / "AppData" / "Roaming"))
    elif IS_MACOS:
        base = Path.home() / "Library" / "Application Support"
    else:
        base = Path(os.environ.get("XDG_CONFIG_HOME", Path.home() / ".config"))
    return base / APP_NAME / "config.json"


CONFIG_PATH = get_config_path()
LEGACY_CONFIG_PATH = Path.home() / ".pptxtranslator_config.json"
HISTORY_PATH = CONFIG_PATH.parent / "history.json"
CHAT_SESSIONS_DIR = CONFIG_PATH.parent / "chat_sessions"


def enable_windows_dpi_awareness():
    """Keep text and controls crisp on scaled Windows displays."""
    if not IS_WINDOWS:
        return
    try:
        import ctypes
        ctypes.windll.shcore.SetProcessDpiAwareness(1)
    except (AttributeError, OSError):
        try:
            ctypes.windll.user32.SetProcessDPIAware()
        except (AttributeError, OSError):
            pass


def set_window_icon(window):
    """Set a supported runtime icon without assuming Windows paths."""
    png_path = resource_path("pptx_icon.png")
    ico_path = resource_path("pptx_icon.ico")
    try:
        if IS_WINDOWS and ico_path.exists():
            window.iconbitmap(str(ico_path))
        elif png_path.exists():
            icon = tk.PhotoImage(file=str(png_path))
            window.iconphoto(True, icon)
            window._app_icon = icon
    except (tk.TclError, OSError):
        pass


def center_window(window, width, height, parent=None):
    """Size and center a window while keeping it on the visible display."""
    window.update_idletasks()
    if parent and parent.winfo_viewable():
        x = parent.winfo_rootx() + max(0, (parent.winfo_width() - width) // 2)
        y = parent.winfo_rooty() + max(0, (parent.winfo_height() - height) // 2)
    else:
        x = max(0, (window.winfo_screenwidth() - width) // 2)
        y = max(0, (window.winfo_screenheight() - height) // 2)
    window.geometry(f"{width}x{height}+{x}+{y}")


def open_path(path):
    """Open a file or directory with the OS default application."""
    target = str(Path(path).resolve())
    if IS_WINDOWS:
        os.startfile(target)  # type: ignore[attr-defined]
    elif IS_MACOS:
        subprocess.Popen(["open", target])
    else:
        subprocess.Popen(["xdg-open", target])


def reveal_path(path):
    """Reveal a file in Finder/Explorer, or open its directory elsewhere."""
    target = Path(path).resolve()
    if IS_WINDOWS:
        subprocess.Popen(["explorer", "/select,", str(target)])
    elif IS_MACOS:
        subprocess.Popen(["open", "-R", str(target)])
    else:
        open_path(target if target.is_dir() else target.parent)


if TkinterDnD is not None:
    class DragDropCTk(ctk.CTk, TkinterDnD.DnDWrapper):
        """CustomTkinter root with native Finder/Explorer file drops."""

        def __init__(self, *args, **kwargs):
            ctk.CTk.__init__(self, *args, **kwargs)
            self.TkdndVersion = TkinterDnD._require(self)
else:
    DragDropCTk = ctk.CTk


class Tooltip:
    """Small delayed tooltip for icon-only controls."""

    def __init__(self, widget, text, delay=500):
        self.widget = widget
        self.text = text
        self.delay = delay
        self._job = None
        self._window = None
        widget.bind("<Enter>", self._schedule, add="+")
        widget.bind("<Leave>", self._hide, add="+")
        widget.bind("<ButtonPress>", self._hide, add="+")

    def _schedule(self, _event=None):
        self._cancel()
        self._job = self.widget.after(self.delay, self._show)

    def _cancel(self):
        if self._job:
            self.widget.after_cancel(self._job)
            self._job = None

    def _show(self):
        if self._window or not self.widget.winfo_exists():
            return
        x = self.widget.winfo_rootx() + self.widget.winfo_width() // 2
        y = self.widget.winfo_rooty() + self.widget.winfo_height() + 8
        self._window = ctk.CTkToplevel(self.widget)
        self._window.withdraw()
        self._window.overrideredirect(True)
        self._window.attributes("-topmost", True)
        label = ctk.CTkLabel(
            self._window, text=self.text, height=26, corner_radius=6,
            fg_color=("#242832", "#F4F6FA"),
            text_color=("#FFFFFF", "#171A21"), font=ui_font(10),
        )
        label.pack()
        self._window.update_idletasks()
        x -= self._window.winfo_width() // 2
        self._window.geometry(f"+{x}+{y}")
        self._window.deiconify()

    def _hide(self, _event=None):
        self._cancel()
        if self._window:
            self._window.destroy()
            self._window = None


class ModernComboBox(ctk.CTkFrame):
    """Fully clickable selector with a compact, searchable animated popup."""

    def __init__(self, master, variable=None, values=(), command=None,
                 width=220, height=40, max_visible=7, searchable=None, **kwargs):
        super().__init__(
            master, width=width, height=height, corner_radius=7,
            fg_color=COLORS["surface"], border_width=1,
            border_color=COLORS["border"], **kwargs,
        )
        self._values = list(values)
        self._command = command
        self._state = "normal"
        self._popup = None
        self._outside_binding = None
        self._max_visible = max(3, max_visible)
        self._searchable = len(self._values) > 10 if searchable is None else searchable
        self._variable = variable or ctk.StringVar(
            master=self, value=self._values[0] if self._values else ""
        )
        self._trace_id = self._variable.trace_add("write", self._on_value_changed)

        self.grid_propagate(False)
        self.grid_columnconfigure(0, weight=1)
        self._label = ctk.CTkLabel(
            self, textvariable=self._variable, anchor="w",
            text_color=COLORS["text"], font=ui_font(12),
        )
        self._label.grid(row=0, column=0, sticky="nsew", padx=(12, 4), pady=2)
        self._arrow = ctk.CTkLabel(
            self, text="", image=load_icon("chevron-down", 15),
            width=30, height=30,
        )
        self._arrow.grid(row=0, column=1, padx=(2, 5), pady=4)

        for widget in (self, self._label, self._arrow):
            widget.bind("<Button-1>", self._toggle_popup, add="+")
            widget.bind("<Enter>", self._on_enter, add="+")
            widget.bind("<Leave>", self._on_leave, add="+")

    def get(self):
        return self._variable.get()

    def set(self, value):
        self._variable.set(value)

    def animate_content(self, offset=0, tone="normal"):
        """Apply a temporary directional/tint treatment during value changes."""
        self._label.grid_configure(padx=(max(5, 12 + offset), 4))
        tones = {
            "muted": COLORS["muted"],
            "accent": COLORS["accent"],
            "normal": COLORS["text"],
        }
        self._label.configure(text_color=tones.get(tone, COLORS["text"]))
        self.configure(
            border_color=COLORS["accent"] if tone != "normal" else COLORS["border"]
        )

    def configure(self, require_redraw=False, **kwargs):
        state = kwargs.pop("state", None)
        values = kwargs.pop("values", None)
        command = kwargs.pop("command", None)
        result = super().configure(require_redraw=require_redraw, **kwargs)
        if values is not None and hasattr(self, "_values"):
            self._values = list(values)
        if command is not None and hasattr(self, "_command"):
            self._command = command
        if state is not None and hasattr(self, "_state"):
            self._state = state
            disabled = state == "disabled"
            self._label.configure(
                text_color=COLORS["muted"] if disabled else COLORS["text"]
            )
            if disabled:
                self.close_popup()
        return result

    def _on_value_changed(self, *_args):
        if self._popup and self._popup.winfo_exists():
            self._render_options()

    def _on_enter(self, _event=None):
        if self._state == "normal" and not self._popup:
            self.configure(border_color=COLORS["muted"])

    def _on_leave(self, _event=None):
        if self._popup:
            return
        x, y = self.winfo_pointerxy()
        under_pointer = self.winfo_containing(x, y)
        if under_pointer and str(under_pointer).startswith(str(self)):
            return
        self.configure(border_color=COLORS["border"])

    def _toggle_popup(self, _event=None):
        if self._state == "disabled":
            return "break"
        if self._popup:
            self.close_popup()
        else:
            self.open_popup()
        return "break"

    def open_popup(self):
        if self._popup or not self._values:
            return
        self.update_idletasks()
        owner = self.winfo_toplevel()
        popup = ctk.CTkToplevel(owner)
        popup.withdraw()
        popup.overrideredirect(True)
        popup.transient(owner)
        popup.configure(fg_color=COLORS["surface"])
        try:
            popup.attributes("-topmost", True)
            popup.attributes("-alpha", 0.0)
        except tk.TclError:
            pass
        self._popup = popup
        self.configure(border_color=COLORS["accent"])

        container = ctk.CTkFrame(
            popup, corner_radius=8, fg_color=COLORS["surface"],
            border_width=1, border_color=COLORS["border"],
        )
        container.pack(fill="both", expand=True)

        if self._searchable:
            self._search_var = ctk.StringVar()
            search = ctk.CTkEntry(
                container, textvariable=self._search_var,
                placeholder_text="Search languages...", height=34,
                corner_radius=6, border_width=1, border_color=COLORS["border"],
                fg_color=COLORS["surface_alt"], font=ui_font(11),
            )
            search.pack(fill="x", padx=8, pady=(8, 4))
            self._search_trace = self._search_var.trace_add(
                "write", lambda *_args: self._render_options()
            )
        else:
            self._search_var = None
            self._search_trace = None

        list_height = self._max_visible * 36
        needs_scroll = self._searchable or len(self._values) > self._max_visible
        if needs_scroll:
            self._options_frame = ctk.CTkScrollableFrame(
                container, height=list_height, corner_radius=0,
                fg_color="transparent", scrollbar_button_color=COLORS["border"],
                scrollbar_button_hover_color=COLORS["muted"],
            )
        else:
            self._options_frame = ctk.CTkFrame(
                container, height=list_height, corner_radius=0,
                fg_color="transparent",
            )
        self._options_frame.pack(fill="both", expand=True, padx=4, pady=(2, 6))
        self._render_options()

        popup_width = max(self.winfo_width(), 220)
        visible = min(self._max_visible, max(1, len(self._values)))
        popup_height = visible * 36 + (52 if self._searchable else 10)
        x = self.winfo_rootx()
        target_y = self.winfo_rooty() + self.winfo_height() + 6
        screen_bottom = self.winfo_screenheight() - 16
        if target_y + popup_height > screen_bottom:
            target_y = self.winfo_rooty() - popup_height - 6
        self._popup_geometry = (popup_width, popup_height, x, target_y)
        popup.geometry(
            f"{popup_width}x{popup_height}+{x}+{target_y - 8}"
        )
        popup.deiconify()
        popup.lift()
        popup.bind("<Escape>", lambda _event: self.close_popup())
        popup.bind("<FocusOut>", self._popup_focus_out, add="+")
        self._outside_binding = owner.bind(
            "<Button-1>", self._owner_clicked, add="+"
        )
        if self._searchable:
            popup.after(90, search.focus_set)
        else:
            popup.focus_set()
        self._animate_popup_open(0)

    def _animate_popup_open(self, step):
        popup = self._popup
        if not popup or not popup.winfo_exists():
            return
        width, height, x, target_y = self._popup_geometry
        progress = min(1.0, step / 7)
        eased = 1 - (1 - progress) ** 3
        y = int(target_y - 8 + 8 * eased)
        popup.geometry(f"{width}x{height}+{x}+{y}")
        try:
            popup.attributes("-alpha", 0.2 + 0.8 * eased)
        except tk.TclError:
            pass
        if step < 7:
            popup.after(14, lambda: self._animate_popup_open(step + 1))

    def _render_options(self):
        if not self._popup or not self._popup.winfo_exists():
            return
        for child in self._options_frame.winfo_children():
            child.destroy()
        query = self._search_var.get().strip().casefold() if self._search_var else ""
        filtered = [value for value in self._values if query in value.casefold()]
        if not filtered:
            ctk.CTkLabel(
                self._options_frame, text="No matching languages",
                text_color=COLORS["muted"], font=ui_font(10), height=42,
            ).pack(fill="x", padx=4, pady=4)
            return
        current = self._variable.get()
        for value in filtered:
            selected = value == current
            option = ctk.CTkButton(
                self._options_frame, text=value, anchor="w", height=34,
                corner_radius=6, fg_color=(
                    ("#EEF0FF", "#292E52") if selected else "transparent"
                ),
                hover_color=COLORS["surface_alt"],
                text_color=COLORS["accent"] if selected else COLORS["text"],
                font=ui_font(11, "bold" if selected else "normal"),
                command=lambda item=value: self._select(item),
            )
            option.pack(fill="x", padx=3, pady=1)

    def _select(self, value):
        changed = value != self._variable.get()
        self._variable.set(value)
        self.close_popup()
        if changed and self._command:
            self._command(value)
        if changed:
            self.animate_content(tone="accent")
            self.after(140, lambda: self.animate_content(tone="normal"))

    def _owner_clicked(self, event):
        widget_path = str(event.widget)
        if widget_path.startswith(str(self)):
            return
        self.close_popup()

    def _popup_focus_out(self, _event=None):
        self.after(80, self._close_if_focus_left)

    def _close_if_focus_left(self):
        if not self._popup or not self._popup.winfo_exists():
            return
        focus = self._popup.focus_get()
        if focus is None:
            self.close_popup()
            return
        focus_path = str(focus)
        owner_path = str(self.winfo_toplevel())
        if not (
            focus_path.startswith(str(self._popup))
            or focus_path.startswith(owner_path)
        ):
            self.close_popup()

    def close_popup(self):
        popup = self._popup
        if not popup:
            return
        self._popup = None
        owner = self.winfo_toplevel()
        if self._outside_binding:
            owner.unbind("<Button-1>", self._outside_binding)
            self._outside_binding = None
        if popup.winfo_exists():
            popup.destroy()
        self.configure(border_color=COLORS["border"])

    def destroy(self):
        self.close_popup()
        if hasattr(self, "_trace_id"):
            self._variable.trace_remove("write", self._trace_id)
        super().destroy()

MODELS = {
    "gpt-4.1-nano": {"label": "4.1 Nano", "input_cost": 0.10, "output_cost": 0.40, "vision": False},
    "gpt-4.1-mini": {"label": "4.1 Mini", "input_cost": 0.40, "output_cost": 1.60, "vision": True},
    "gpt-4.1": {"label": "4.1", "input_cost": 2.00, "output_cost": 8.00, "vision": True},
    "gpt-4o-mini": {"label": "4o Mini", "input_cost": 0.15, "output_cost": 0.60, "vision": True},
    "gpt-4o": {"label": "4o", "input_cost": 2.50, "output_cost": 10.00, "vision": True},
    "gpt-5.4": {"label": "5.4", "input_cost": 5.00, "output_cost": 20.00, "vision": True},
}

AUTO_DETECT = "Auto-detect"

# Keep these as human-readable names: they are shown in the UI and passed directly
# to the translation model, avoiding a separate locale-code mapping layer.
LANGUAGES = (
    "Arabic",
    "Bengali",
    "Bulgarian",
    "Catalan",
    "Chinese (Simplified)",
    "Chinese (Traditional)",
    "Croatian",
    "Czech",
    "Danish",
    "Dutch",
    "English",
    "Estonian",
    "Filipino",
    "Finnish",
    "French",
    "German",
    "Greek",
    "Hebrew",
    "Hindi",
    "Hungarian",
    "Indonesian",
    "Italian",
    "Japanese",
    "Korean",
    "Latvian",
    "Lithuanian",
    "Malay",
    "Norwegian",
    "Persian",
    "Polish",
    "Portuguese (Brazil)",
    "Portuguese (Portugal)",
    "Romanian",
    "Russian",
    "Serbian",
    "Slovak",
    "Slovenian",
    "Spanish",
    "Swedish",
    "Tamil",
    "Thai",
    "Turkish",
    "Ukrainian",
    "Urdu",
    "Vietnamese",
)

TARGET_BATCH_CHARS = 2000
SEP_TOKEN = "\n[SEP]\n"


def load_config():
    config_path = CONFIG_PATH if CONFIG_PATH.exists() else LEGACY_CONFIG_PATH
    try:
        with config_path.open("r", encoding="utf-8") as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError, OSError):
        return {}


def save_config(config):
    try:
        CONFIG_PATH.parent.mkdir(parents=True, exist_ok=True)
        with CONFIG_PATH.open("w", encoding="utf-8") as f:
            json.dump(config, f, indent=2)
    except OSError:
        # Preferences are convenient, but should never stop a translation.
        pass


def _load_json(path, fallback):
    try:
        with Path(path).open("r", encoding="utf-8") as handle:
            return json.load(handle)
    except (FileNotFoundError, json.JSONDecodeError, OSError):
        return fallback


def _save_json(path, payload):
    """Write persistent app data atomically so interrupted writes stay recoverable."""
    target = Path(path)
    temp_path = target.with_suffix(target.suffix + ".tmp")
    try:
        target.parent.mkdir(parents=True, exist_ok=True)
        with temp_path.open("w", encoding="utf-8") as handle:
            json.dump(payload, handle, indent=2, ensure_ascii=False)
        temp_path.replace(target)
        return True
    except OSError:
        try:
            temp_path.unlink(missing_ok=True)
        except OSError:
            pass
        return False


def load_translation_history():
    history = _load_json(HISTORY_PATH, [])
    return history if isinstance(history, list) else []


def save_translation_history(history):
    return _save_json(HISTORY_PATH, history[:200])


def chat_session_path(file_path):
    resolved = Path(file_path).expanduser().resolve()
    normalized = os.path.normcase(str(resolved))
    try:
        stat = resolved.stat()
        identity = f"{normalized}:{stat.st_size}:{stat.st_mtime_ns}"
    except OSError:
        identity = normalized
    digest = hashlib.sha256(identity.encode("utf-8")).hexdigest()
    return CHAT_SESSIONS_DIR / f"{digest}.json"


def load_chat_session(file_path):
    payload = _load_json(chat_session_path(file_path), {})
    messages = payload.get("messages", []) if isinstance(payload, dict) else []
    return [
        {"role": item["role"], "content": item["content"]}
        for item in messages
        if isinstance(item, dict)
        and item.get("role") in {"user", "assistant"}
        and isinstance(item.get("content"), str)
    ]


def save_chat_session(file_path, messages):
    serializable = [
        {"role": item["role"], "content": item["content"]}
        for item in messages
        if item.get("role") in {"user", "assistant"}
        and isinstance(item.get("content"), str)
    ]
    return _save_json(
        chat_session_path(file_path),
        {
            "file_path": str(Path(file_path).resolve()),
            "updated_at": datetime.now(timezone.utc).isoformat(),
            "messages": serializable,
        },
    )


def format_history_time(iso_value):
    try:
        timestamp = datetime.fromisoformat(iso_value.replace("Z", "+00:00"))
        return timestamp.astimezone().strftime("%b %d, %Y  %H:%M")
    except (AttributeError, TypeError, ValueError):
        return "Unknown date"


def soft_wrap_path(path):
    """Add invisible break opportunities without changing the visible path."""
    return str(path).replace("\\", "\\\u200b").replace("/", "/\u200b")


def has_translatable_text(text):
    """Return True for text containing letters from any Unicode script."""
    return bool(text and text.strip()) and any(char.isalpha() for char in text)


def translate_text_batch(texts, model, client, source_language, target_language):
    joined = SEP_TOKEN.join(texts)
    if source_language == AUTO_DETECT:
        source_instruction = (
            "Detect the language used in each segment, then translate it"
        )
    else:
        source_instruction = f"Translate each {source_language} segment"

    messages = [
        {
            "role": "system",
            "content": (
                f"You are a professional document translator. {source_instruction} "
                f"into {target_language}. If a segment is not written in the source "
                "language, leave that segment unchanged. Preserve meaning, tone, line "
                "breaks, numbers, names, and formatting. Do not translate URLs or code. "
                "Each segment is separated by [SEP]. Return the results in the same "
                "order, separated by [SEP], with no commentary or additional text."
            ),
        },
        {"role": "user", "content": joined},
    ]

    try:
        response = client.chat.completions.create(
            model=model, messages=messages, temperature=0.1
        )
        translated_text = response.choices[0].message.content
        input_tokens = response.usage.prompt_tokens
        output_tokens = response.usage.completion_tokens

        translations = [t.strip() for t in translated_text.split("[SEP]")]

        if len(translations) < len(texts):
            translations.extend(texts[len(translations):])
        elif len(translations) > len(texts):
            translations = translations[: len(texts)]

        return translations, input_tokens, output_tokens
    except Exception as e:
        print(f"Error during translation: {e}")
        return texts, 0, 0


def canonicalize_detected_language(response_text):
    """Map a model response back to one of the language selector values."""
    normalized = re.sub(r"[^a-z]+", " ", response_text.lower()).strip()
    aliases = {
        "traditional chinese": "Chinese (Traditional)",
        "chinese traditional": "Chinese (Traditional)",
        "simplified chinese": "Chinese (Simplified)",
        "chinese simplified": "Chinese (Simplified)",
        "brazilian portuguese": "Portuguese (Brazil)",
        "portuguese brazil": "Portuguese (Brazil)",
        "european portuguese": "Portuguese (Portugal)",
        "portuguese portugal": "Portuguese (Portugal)",
        "farsi": "Persian",
        "tagalog": "Filipino",
    }
    for alias, language in aliases.items():
        if normalized == alias or alias in normalized:
            return language

    for language in LANGUAGES:
        candidate = re.sub(r"[^a-z]+", " ", language.lower()).strip()
        if normalized == candidate or candidate in normalized:
            return language
    return None


def detect_document_language(text_sample, model, client):
    """Ask the model for the dominant document language."""
    if not has_translatable_text(text_sample):
        return None, 0, 0

    language_list = ", ".join(LANGUAGES)
    messages = [
        {
            "role": "system",
            "content": (
                "Identify the dominant language in the document sample. Distinguish "
                "Simplified Chinese from Traditional Chinese and Brazilian Portuguese "
                "from European Portuguese. Reply with exactly one language name from "
                f"this list and no other text: {language_list}."
            ),
        },
        {"role": "user", "content": text_sample[:8000]},
    ]
    try:
        response = client.chat.completions.create(
            model=model, messages=messages, temperature=0
        )
        detected = canonicalize_detected_language(response.choices[0].message.content or "")
        return detected, response.usage.prompt_tokens, response.usage.completion_tokens
    except Exception as exc:
        print(f"Error during language detection: {exc}")
        return None, 0, 0


def build_batches(text_items):
    batches = []
    current_batch = []
    current_chars = 0
    for text in text_items:
        if current_chars + len(text) > TARGET_BATCH_CHARS and current_batch:
            batches.append(current_batch)
            current_batch = []
            current_chars = 0
        current_batch.append(text)
        current_chars += len(text)
    if current_batch:
        batches.append(current_batch)
    return batches


def iter_shape_paragraphs(shape):
    """Yield paragraphs from text boxes, table cells, and grouped shapes."""
    if getattr(shape, "has_text_frame", False):
        yield from shape.text_frame.paragraphs

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                yield from cell.text_frame.paragraphs

    if hasattr(shape, "shapes"):
        for child_shape in shape.shapes:
            yield from iter_shape_paragraphs(child_shape)


def iter_pptx_paragraphs(presentation):
    for slide in presentation.slides:
        for shape in slide.shapes:
            yield from iter_shape_paragraphs(shape)


def extract_text_sample(file_path, max_chars=8000):
    """Read enough document text for language detection without loading images."""
    ext = Path(file_path).suffix.lower()
    parts = []
    char_count = 0

    if ext == ".pdf":
        doc = fitz.open(file_path)
        try:
            for page in doc:
                text = page.get_text().strip()
                if has_translatable_text(text):
                    parts.append(text)
                    char_count += len(text)
                if char_count >= max_chars:
                    break
        finally:
            doc.close()
    else:
        presentation = Presentation(file_path)
        for paragraph in iter_pptx_paragraphs(presentation):
            text = "".join(run.text for run in paragraph.runs).strip()
            if has_translatable_text(text):
                parts.append(text)
                char_count += len(text)
            if char_count >= max_chars:
                break

    return "\n".join(parts)[:max_chars]


def scan_pptx_paragraphs(file_path):
    prs = Presentation(file_path)
    return sum(
        1
        for para in iter_pptx_paragraphs(prs)
        if has_translatable_text("".join(run.text for run in para.runs))
    )


def process_pptx(file_path, output_path, model, client, source_language,
                 target_language, progress_callback=None, cancel_event=None,
                 para_offset=0):
    prs = Presentation(file_path)
    text_items = []
    para_map = []

    for para in iter_pptx_paragraphs(prs):
        full_text = "".join(run.text for run in para.runs)
        if has_translatable_text(full_text):
            text_items.append(full_text)
            para_map.append(para)

    total = len(text_items)
    translated_results = []
    input_tokens = 0
    output_tokens = 0
    cancelled = False

    batches = build_batches(text_items)
    paras_done = 0

    # Report initial progress so progress bar shows 0% even for single-batch files
    if progress_callback and total > 0:
        progress_callback(para_offset, input_tokens, output_tokens)

    for batch in batches:
        if cancel_event and cancel_event.is_set():
            cancelled = True
            break

        translations, batch_in, batch_out = translate_text_batch(
            batch, model, client, source_language, target_language
        )
        input_tokens += batch_in
        output_tokens += batch_out
        translated_results.extend(translations)
        paras_done += len(batch)

        if progress_callback:
            progress_callback(para_offset + paras_done, input_tokens, output_tokens)

    if not cancelled:
        for para, translated_text in zip(para_map, translated_results):
            if para.runs:
                para.runs[0].text = translated_text
                for run in para.runs[1:]:
                    run.text = ""
        prs.save(output_path)

    return input_tokens, output_tokens, not cancelled, total


def extract_pptx_content(file_path):
    """Extract text and images from a PPTX file, organized by slide."""
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slide_data = {"number": i, "texts": [], "images": []}
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_data["texts"].append(text)
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                rows_text = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows_text.append(" | ".join(cells))
                if rows_text:
                    slide_data["texts"].append("[Table]\n" + "\n".join(rows_text))
            try:
                img = shape.image
                b64 = base64.b64encode(img.blob).decode()
                slide_data["images"].append({
                    "base64": b64,
                    "content_type": img.content_type or "image/png",
                })
            except (AttributeError, Exception):
                pass
        slides.append(slide_data)
    return slides


def scan_pdf_paragraphs(file_path):
    """Count text spans that contain letters from any language."""
    doc = fitz.open(file_path)
    count = 0
    for page in doc:
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if block["type"] != 0:  # skip non-text blocks
                continue
            for line in block["lines"]:
                for span in line["spans"]:
                    if has_translatable_text(span["text"]):
                        count += 1
    doc.close()
    return count


def process_pdf(file_path, output_path, model, client, source_language,
                target_language, progress_callback=None, cancel_event=None,
                para_offset=0):
    """Translate text in a PDF, replacing it in-place."""
    doc = fitz.open(file_path)
    text_items = []
    span_info = []  # (page_idx, rect, font_size, font_name, color)

    for page_idx, page in enumerate(doc):
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if block["type"] != 0:
                continue
            for line in block["lines"]:
                for span in line["spans"]:
                    if has_translatable_text(span["text"]):
                        text_items.append(span["text"])
                        span_info.append({
                            "page": page_idx,
                            "bbox": fitz.Rect(span["bbox"]),
                            "size": span["size"],
                            "font": span["font"],
                            "color": span["color"],
                        })

    total = len(text_items)
    translated_results = []
    input_tokens = 0
    output_tokens = 0
    cancelled = False

    batches = build_batches(text_items)
    paras_done = 0

    if progress_callback and total > 0:
        progress_callback(para_offset, input_tokens, output_tokens)

    for batch in batches:
        if cancel_event and cancel_event.is_set():
            cancelled = True
            break

        translations, batch_in, batch_out = translate_text_batch(
            batch, model, client, source_language, target_language
        )
        input_tokens += batch_in
        output_tokens += batch_out
        translated_results.extend(translations)
        paras_done += len(batch)

        if progress_callback:
            progress_callback(para_offset + paras_done, input_tokens, output_tokens)

    if not cancelled:
        # Group spans by page for efficient redaction
        pages_to_redact = {}
        for info, translated in zip(span_info, translated_results):
            pg = info["page"]
            if pg not in pages_to_redact:
                pages_to_redact[pg] = []
            pages_to_redact[pg].append((info, translated))

        for pg_idx, items in pages_to_redact.items():
            page = doc[pg_idx]
            # Add redaction annotations for all spans on this page
            for info, _ in items:
                page.add_redact_annot(info["bbox"])
            # Apply all redactions at once (removes original text)
            page.apply_redactions()
            # Insert translated text
            for info, translated in items:
                if not translated:
                    continue
                rect = info["bbox"]
                fontsize = info["size"]
                # Convert int color to RGB tuple
                c = info["color"]
                rgb = ((c >> 16) & 0xFF, (c >> 8) & 0xFF, c & 0xFF)
                color_hex = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                safe_text = html.escape(translated).replace("\n", "<br>")
                css = (
                    "body, div { margin: 0; padding: 0; } "
                    f"div {{ font-family: sans-serif; font-size: {fontsize}pt; "
                    f"line-height: 1.05; color: {color_hex}; }}"
                )
                # Story-backed HTML insertion supplies Unicode fallback fonts and
                # HarfBuzz shaping for CJK, right-to-left, and complex scripts.
                page.insert_htmlbox(
                    rect,
                    f'<div dir="auto">{safe_text}</div>',
                    css=css,
                    scale_low=0,
                    overlay=True,
                )

        doc.ez_save(output_path)

    doc.close()
    return input_tokens, output_tokens, not cancelled, total


def extract_pdf_content(file_path):
    """Extract text and images from a PDF file, organized by page."""
    doc = fitz.open(file_path)
    pages = []
    for i, page in enumerate(doc, 1):
        page_data = {"number": i, "texts": [], "images": []}
        text = page.get_text().strip()
        if text:
            page_data["texts"].append(text)
        # Extract images
        for img_info in page.get_images(full=True):
            try:
                xref = img_info[0]
                img_data = doc.extract_image(xref)
                if img_data:
                    b64 = base64.b64encode(img_data["image"]).decode()
                    ext = img_data.get("ext", "png")
                    content_type = f"image/{ext}" if ext != "jpg" else "image/jpeg"
                    page_data["images"].append({
                        "base64": b64,
                        "content_type": content_type,
                    })
            except Exception:
                pass
        pages.append(page_data)
    doc.close()
    return pages


def extract_file_content(file_path):
    """Dispatch to PPTX or PDF content extractor based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_pdf_content(file_path)
    else:
        return extract_pptx_content(file_path)


def format_file_size(byte_count):
    value = float(byte_count)
    for unit in ("B", "KB", "MB", "GB"):
        if value < 1024 or unit == "GB":
            return f"{value:.1f} {unit}" if unit != "B" else f"{int(value)} B"
        value /= 1024


def get_document_metadata(filepath):
    path = Path(filepath)
    size = format_file_size(path.stat().st_size)
    try:
        if path.suffix.lower() == ".pptx":
            count = len(Presentation(path).slides)
            unit = "slide" if count == 1 else "slides"
        elif path.suffix.lower() == ".pdf":
            document = fitz.open(path)
            try:
                count = len(document)
            finally:
                document.close()
            unit = "page" if count == 1 else "pages"
        else:
            return size
        return f"{count} {unit}  \u00b7  {size}"
    except Exception:
        return size


class FileEntry(ctk.CTkFrame):
    def __init__(self, master, filepath, on_remove, **kwargs):
        super().__init__(master, **kwargs)
        self.filepath = filepath
        self.on_remove = on_remove
        self.metadata = get_document_metadata(filepath)
        self.configure(
            fg_color="transparent", corner_radius=8,
            border_width=0,
        )
        self.grid_columnconfigure(1, weight=1)

        icon_tile = ctk.CTkLabel(
            self, text="", image=load_icon("presentation", 22),
            width=44, height=44, corner_radius=7,
            fg_color=COLORS["surface_alt"],
        )
        icon_tile.grid(row=0, column=0, rowspan=3, padx=(10, 12), pady=10)

        self.label = ctk.CTkLabel(
            self, text=ellipsize_middle(os.path.basename(filepath), 58), anchor="w",
            text_color=COLORS["text"], font=ui_font(13, "bold"), height=20,
        )
        self.label.grid(row=0, column=1, sticky="sew", pady=(9, 0))
        Tooltip(self.label, os.path.basename(filepath))

        self.meta_label = ctk.CTkLabel(
            self, text=self.metadata, anchor="w", text_color=COLORS["muted"],
            font=ui_font(10), height=18,
        )
        self.meta_label.grid(row=1, column=1, sticky="new", pady=(0, 7))

        self.progress_bar = ctk.CTkProgressBar(
            self, height=3, corner_radius=1,
            fg_color=COLORS["border"], progress_color=COLORS["accent"],
        )
        self.progress_bar.grid(row=2, column=1, sticky="ew", pady=(0, 8))
        self.progress_bar.set(0)
        self.progress_bar.grid_remove()

        self.state_label = ctk.CTkLabel(
            self, text="", width=92, anchor="e", text_color=COLORS["muted"],
            font=ui_font(10, "bold"), compound="left",
        )
        self.state_label.grid(row=0, column=2, rowspan=3, padx=(12, 4), pady=10)

        self.more_btn = ctk.CTkButton(
            self, text="", image=load_icon("ellipsis", 18),
            width=32, height=32, corner_radius=6,
            fg_color="transparent", hover_color=COLORS["border"],
            command=self.show_menu,
        )
        self.more_btn.grid(row=0, column=3, rowspan=3, padx=(2, 8), pady=10)
        Tooltip(self.more_btn, "More actions")

        for widget in (self, self.label, self.meta_label):
            widget.bind("<Enter>", self._on_enter, add="+")
            widget.bind("<Leave>", self._on_leave, add="+")

    def _on_enter(self, _event=None):
        self.configure(fg_color=COLORS["surface_alt"])

    def _on_leave(self, _event=None):
        x, y = self.winfo_pointerxy()
        under_pointer = self.winfo_containing(x, y)
        if under_pointer and str(under_pointer).startswith(str(self)):
            return
        self.configure(fg_color="transparent")

    def show_menu(self):
        menu = tk.Menu(self, tearoff=0)
        menu.add_command(label="Open", command=lambda: open_path(self.filepath))
        menu.add_command(label="Show in folder", command=lambda: reveal_path(self.filepath))
        menu.add_separator()
        menu.add_command(label="Remove", command=self.on_remove)
        try:
            menu.tk_popup(
                self.more_btn.winfo_rootx(),
                self.more_btn.winfo_rooty() + self.more_btn.winfo_height(),
            )
        finally:
            menu.grab_release()

    def set_state(self, state, progress=None, error=None):
        self.meta_label.configure(text=self.metadata, text_color=COLORS["muted"])
        self.state_label.configure(image=None)
        self.progress_bar.grid_remove()

        if state == "Queued":
            self.state_label.configure(text="Queued", text_color=COLORS["muted"])
        elif state == "Translating":
            percent = int((progress or 0) * 100)
            self.state_label.configure(text=f"{percent}%", text_color=COLORS["accent"])
            self.progress_bar.set(progress or 0)
            self.progress_bar.grid()
        elif state == "Complete":
            self.state_label.configure(
                text=" Complete", image=load_icon("check-circle", 14),
                text_color=COLORS["success"],
            )
        elif state == "Failed":
            self.state_label.configure(
                text=" Failed", image=load_icon("circle-alert", 14),
                text_color=COLORS["danger"],
            )
            self.meta_label.configure(
                text=ellipsize_middle(error or "Translation failed", 64),
                text_color=COLORS["danger"],
            )
        else:
            self.state_label.configure(text="Ready", text_color=COLORS["muted"])


class SettingsPanel:
    """Centered in-app settings overlay without a separate OS window."""

    def __init__(self, app, focus_api=False):
        self.app = app
        self.show_key = False
        self._closed = False
        self._closing = False

        self.backdrop = ctk.CTkFrame(
            app.master, corner_radius=0,
            fg_color=("#E9EDF4", "#090B0E"),
        )
        self.backdrop.place(x=0, y=0, relwidth=1, relheight=1)
        self.backdrop.lift()

        self.card = ctk.CTkFrame(
            self.backdrop, width=500, height=460, corner_radius=12,
            fg_color=COLORS["surface"], border_width=1,
            border_color=COLORS["border"],
        )
        self.card.place(relx=0.5, rely=0.525, anchor="center")
        self.card.grid_propagate(False)
        self.card.grid_columnconfigure(0, weight=1)
        self.card.grid_rowconfigure(1, weight=1)
        self._build_ui()

        self._escape_binding = app.master.bind(
            "<Escape>", lambda _event: self.close(), add="+"
        )
        self._animate_open(0)
        if focus_api:
            self.backdrop.after(170, self.api_key_entry.focus_set)

    @property
    def is_open(self):
        return not self._closed and self.backdrop.winfo_exists()

    def _build_ui(self):
        header = ctk.CTkFrame(
            self.card, height=60, corner_radius=11, fg_color=COLORS["surface"]
        )
        header.grid(row=0, column=0, sticky="ew")
        header.grid_propagate(False)
        header.grid_columnconfigure(0, weight=1)
        ctk.CTkLabel(
            header, text="Settings", anchor="w", text_color=COLORS["text"],
            font=ui_font(17, "bold"),
        ).grid(row=0, column=0, sticky="w", padx=22, pady=17)
        close_btn = ctk.CTkButton(
            header, text="", image=load_icon("x", 18), width=32, height=32,
            corner_radius=6, fg_color="transparent", hover_color=COLORS["border"],
            command=self.close,
        )
        close_btn.grid(row=0, column=1, padx=16, pady=14)
        Tooltip(close_btn, "Close")

        body = ctk.CTkFrame(self.card, fg_color="transparent")
        body.grid(row=1, column=0, sticky="nsew", padx=30, pady=(24, 20))
        body.grid_columnconfigure(0, weight=1)

        ctk.CTkLabel(
            body, text="OpenAI API key", anchor="w", text_color=COLORS["text"],
            font=ui_font(11, "bold"),
        ).grid(row=0, column=0, sticky="ew", pady=(0, 7))
        key_shell = ctk.CTkFrame(
            body, height=42, fg_color=COLORS["surface"], corner_radius=7,
            border_width=1, border_color=COLORS["border"],
        )
        key_shell.grid(row=1, column=0, sticky="ew")
        key_shell.grid_columnconfigure(0, weight=1)
        self.api_key_entry = ctk.CTkEntry(
            key_shell, textvariable=self.app.api_key_var, show="\u2022",
            placeholder_text="sk-...", height=38, border_width=0,
            fg_color="transparent", font=ui_font(12),
        )
        self.api_key_entry.grid(row=0, column=0, sticky="ew", padx=(8, 0), pady=1)
        self.api_key_entry.bind("<FocusOut>", lambda _event: self.app.save_api_key())
        self.key_btn = ctk.CTkButton(
            key_shell, text="", image=load_icon("eye", 17), width=34, height=32,
            corner_radius=5, fg_color="transparent", hover_color=COLORS["border"],
            command=self.toggle_key_visibility,
        )
        self.key_btn.grid(row=0, column=1, padx=(2, 4), pady=4)
        self.key_tooltip = Tooltip(self.key_btn, "Show API key")
        ctk.CTkLabel(
            body, text="Stored only in your local app settings.", anchor="w",
            text_color=COLORS["muted"], font=ui_font(9),
        ).grid(row=2, column=0, sticky="ew", pady=(5, 18))

        ctk.CTkLabel(
            body, text="Translation model", anchor="w", text_color=COLORS["text"],
            font=ui_font(11, "bold"),
        ).grid(row=3, column=0, sticky="ew", pady=(0, 7))
        self.model_menu = ModernComboBox(
            body, variable=self.app.model_var,
            values=[MODELS[mid]["label"] for mid in self.app.model_ids],
            height=42, max_visible=6, searchable=False,
            command=self.app.on_model_change,
        )
        self.model_menu.grid(row=4, column=0, sticky="ew", pady=(0, 18))

        ctk.CTkLabel(
            body, text="Appearance", anchor="w", text_color=COLORS["text"],
            font=ui_font(11, "bold"),
        ).grid(row=5, column=0, sticky="ew", pady=(0, 7))
        self.appearance_menu = ModernComboBox(
            body, variable=self.app.appearance_var,
            values=["System", "Light", "Dark"], height=42,
            max_visible=3, searchable=False,
            command=self._queue_appearance_change,
        )
        self.appearance_menu.grid(row=6, column=0, sticky="ew")

        ctk.CTkButton(
            self.card, text="Done", width=92, height=38, corner_radius=7,
            fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
            font=ui_font(11, "bold"), command=self.close,
        ).grid(row=2, column=0, sticky="e", padx=28, pady=(0, 22))

    def _animate_open(self, step):
        if self._closed or not self.backdrop.winfo_exists():
            return
        progress = min(1.0, step / 8)
        eased = 1 - (1 - progress) ** 3
        self.card.place_configure(rely=0.5 + 0.025 * (1 - eased))
        if step < 8:
            self.backdrop.after(14, lambda: self._animate_open(step + 1))

    def _queue_appearance_change(self, appearance):
        # Let the popup finish destroying its native Tk resources before the
        # global CustomTkinter appearance tracker redraws every widget.
        self.appearance_menu.close_popup()
        self.backdrop.after(
            25, lambda value=appearance: self.app.change_appearance_mode(value)
        )

    def toggle_key_visibility(self):
        self.show_key = not self.show_key
        self.api_key_entry.configure(show="" if self.show_key else "\u2022")
        self.key_btn.configure(
            image=load_icon("eye-off" if self.show_key else "eye", 17)
        )
        self.key_tooltip.text = "Hide API key" if self.show_key else "Show API key"

    def close(self):
        if self._closed or self._closing:
            return
        self._closing = True
        self.model_menu.close_popup()
        self.appearance_menu.close_popup()
        self.app.save_api_key()
        if self._escape_binding:
            self.app.master.unbind("<Escape>", self._escape_binding)
            self._escape_binding = None
        self._animate_close(0)

    def _animate_close(self, step):
        if self._closed or not self.backdrop.winfo_exists():
            return
        progress = min(1.0, step / 8)
        eased = progress ** 2
        self.card.place_configure(rely=0.5 + 0.025 * eased)
        if step < 8:
            self.backdrop.after(14, lambda: self._animate_close(step + 1))
            return
        self._closed = True
        self.app.settings_window = None
        self.backdrop.destroy()
        self.app.master.focus_set()


class ChatPanel:
    """Ask AI workspace embedded in the main application window."""

    def __init__(self, app, pptx_path, model_id, api_key, on_close=None):
        self.app = app
        self.master = app.master
        self.win = app.master  # Existing async callbacks use the root event loop.
        self.pptx_path = pptx_path
        self.model_id = model_id
        self.api_key = api_key
        self.client = OpenAI(api_key=api_key)
        self.messages = []
        self.saved_messages = load_chat_session(pptx_path)
        self.streaming = False
        self._closed = False
        self.on_close = on_close
        self._current_ai_label = None

        self.overlay = ctk.CTkFrame(
            self.master, corner_radius=0, fg_color=COLORS["background"]
        )
        self.overlay.place(x=0, y=0, relwidth=1, relheight=1)
        self.overlay.lift()

        self.setup_ui()
        self.load_presentation()

    def setup_ui(self):
        # Header
        header = ctk.CTkFrame(
            self.overlay, height=58, corner_radius=0, fg_color=COLORS["surface"],
            border_width=0,
        )
        header.pack(fill="x")
        header.pack_propagate(False)
        header.grid_columnconfigure(2, weight=1)

        ctk.CTkButton(
            header, text="Back", width=62, height=32, corner_radius=6,
            fg_color="transparent", hover_color=COLORS["surface_alt"],
            text_color=COLORS["text"], font=ui_font(10, "bold"),
            command=self.close,
        ).grid(row=0, column=0, padx=(14, 4), pady=13)

        ctk.CTkLabel(
            header, text="Ask AI", text_color=COLORS["text"],
            font=ui_font(16, "bold"),
        ).grid(row=0, column=1, padx=(8, 12), pady=12)

        ctk.CTkLabel(
            header, text=ellipsize_middle(os.path.basename(self.pptx_path), 38),
            font=ui_font(12),
            text_color=COLORS["muted"], anchor="w",
        ).grid(row=0, column=2, sticky="w", padx=(0, 16), pady=10)

        model_label = MODELS.get(self.model_id, {}).get("label", self.model_id)
        ctk.CTkLabel(
            header, text=f"Model: {model_label}",
            font=ui_font(11),
            text_color=COLORS["muted"],
        ).grid(row=0, column=3, padx=(8, 12), pady=10)
        ctk.CTkButton(
            header, text="Open File", image=load_icon("presentation", 15),
            compound="left", width=104, height=34, corner_radius=7,
            fg_color="transparent", hover_color=COLORS["surface_alt"],
            border_width=1, border_color=COLORS["border"],
            text_color=COLORS["text"], font=ui_font(10, "bold"),
            command=self.open_document,
        ).grid(row=0, column=4, padx=(0, 16), pady=12)

        # Chat area
        self.chat_frame = ctk.CTkScrollableFrame(
            self.overlay, corner_radius=0, fg_color=COLORS["background"]
        )
        self.chat_frame.pack(fill="both", expand=True, padx=0, pady=0)

        # Input area
        input_frame = ctk.CTkFrame(
            self.overlay, corner_radius=0, height=70, fg_color=COLORS["surface"]
        )
        input_frame.pack(fill="x", side="bottom")
        input_frame.pack_propagate(False)

        self.input_var = ctk.StringVar()
        self.input_entry = ctk.CTkEntry(
            input_frame, textvariable=self.input_var,
            placeholder_text="Ask about the presentation...",
            height=44, corner_radius=10, border_color=COLORS["border"],
            font=ui_font(13),
        )
        self.input_entry.pack(side="left", fill="x", expand=True, padx=(12, 8), pady=11)
        self.input_entry.bind("<Return>", lambda e: self.send_message())

        self.send_btn = ctk.CTkButton(
            input_frame, text="Send", width=82, height=44, corner_radius=10,
            fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
            font=ui_font(13, "bold"),
            command=self.send_message,
        )
        self.send_btn.pack(side="right", padx=(0, 12), pady=11)

    def open_document(self):
        try:
            open_path(self.pptx_path)
        except OSError as exc:
            messagebox.showerror(
                "Unable to open file", str(exc), parent=self.win
            )

    def close(self):
        if self._closed:
            return
        self._closed = True
        self._save_session()
        self.app.chat_panel = None
        if self.overlay.winfo_exists():
            self.overlay.destroy()
        if self.on_close:
            self.on_close()

    def _save_session(self):
        save_chat_session(self.pptx_path, self.messages)

    def load_presentation(self):
        self.add_system_bubble("Loading presentation...")
        threading.Thread(target=self._load_pptx, daemon=True).start()

    def _load_pptx(self):
        try:
            self.slides_content = extract_file_content(self.pptx_path)

            # Build system message with text content
            text_summary = ""
            for slide in self.slides_content:
                text_summary += f"\n--- Slide {slide['number']} ---\n"
                for t in slide["texts"]:
                    text_summary += t + "\n"
                if slide["images"]:
                    text_summary += f"[{len(slide['images'])} image(s) on this slide]\n"

            system_msg = (
                "You are a helpful AI assistant analyzing a PowerPoint presentation. "
                "Here is the text content of each slide:\n" + text_summary + "\n"
                "Answer the user's questions about this presentation. Be specific about "
                "slide numbers when relevant. If the user asks about visual content, "
                "describe what you can see in the images provided."
            )
            self.messages = [{"role": "system", "content": system_msg}]

            # Include images for vision-capable models
            model_info = MODELS.get(self.model_id, {})
            has_vision = model_info.get("vision", False)

            if has_vision:
                content_parts = [{"type": "text", "text": "Here are the images from the presentation for your reference. Please acknowledge."}]
                img_count = 0
                for slide in self.slides_content:
                    for img in slide["images"]:
                        content_parts.append({
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{img['content_type']};base64,{img['base64']}",
                                "detail": "low",
                            }
                        })
                        img_count += 1
                if img_count > 0:
                    self.messages.append({"role": "user", "content": content_parts})
                    self.messages.append({
                        "role": "assistant",
                        "content": f"I've received and analyzed all {img_count} image(s) from the presentation. I'm ready to answer your questions about both the text and visual content."
                    })

            self.messages.extend(self.saved_messages)

            total_slides = len(self.slides_content)
            total_images = sum(len(s["images"]) for s in self.slides_content)

            if not self._closed:
                self.win.after(0, lambda: self._on_pptx_loaded(total_slides, total_images))
        except Exception as e:
            if not self._closed:
                self.win.after(0, lambda err=str(e): self._on_pptx_error(err))

    def _on_pptx_loaded(self, num_slides, num_images):
        if self._closed:
            return
        for w in self.chat_frame.winfo_children():
            w.destroy()

        if self.saved_messages:
            self.add_system_bubble(
                f"Restored {len(self.saved_messages)} saved message(s) for this file."
            )
            for message in self.saved_messages:
                if message["role"] == "user":
                    self.add_user_bubble(message["content"])
                else:
                    self.add_ai_bubble(message["content"])
            self.input_entry.focus()
            return

        img_text = f" and {num_images} image(s)" if num_images > 0 else ""
        self.add_ai_bubble(
            f"I've analyzed your presentation ({num_slides} slide(s){img_text}). "
            "Ask me anything! For example:\n\n"
            "\u2022  \"What is slide 3 about?\"\n"
            "\u2022  \"Summarize the entire presentation\"\n"
            "\u2022  \"What are the key points?\"\n"
            "\u2022  \"Explain the diagram on slide 5\""
        )
        self.input_entry.focus()

    def _on_pptx_error(self, error):
        if self._closed:
            return
        for w in self.chat_frame.winfo_children():
            w.destroy()
        self.add_system_bubble(f"Error loading presentation: {error}")

    def add_user_bubble(self, text):
        outer = ctk.CTkFrame(self.chat_frame, fg_color="transparent")
        outer.pack(fill="x", padx=12, pady=(6, 2))

        inner = ctk.CTkFrame(outer, fg_color=COLORS["accent"], corner_radius=16)
        inner.pack(side="right")

        label = ctk.CTkLabel(
            inner, text=text, wraplength=420, justify="left",
            text_color="white", font=ui_font(13),
        )
        label.pack(padx=14, pady=10)

        self._scroll_to_bottom()

    def add_ai_bubble(self, text="", formatted=True):
        outer = ctk.CTkFrame(self.chat_frame, fg_color="transparent")
        outer.pack(fill="x", padx=12, pady=(6, 2))

        bubble_color = ("#FFFFFF", "#22262D")
        inner = ctk.CTkFrame(outer, fg_color=bubble_color, corner_radius=12)
        inner.pack(side="left", anchor="nw")

        textbox = tk.Text(
            inner, wrap="word", borderwidth=0, highlightthickness=0,
            bg=resolve_color(bubble_color), fg=resolve_color(COLORS["text"]),
            font=(UI_FONT_FAMILY, 11),
            cursor="arrow", padx=14, pady=10,
            selectbackground=resolve_color(COLORS["accent"]), selectforeground="white",
            width=52, height=1, relief="flat",
            insertbackground=resolve_color(COLORS["text"]),
        )
        textbox.pack(fill="x")

        self._configure_markdown_tags(textbox)

        if text:
            if formatted:
                self._insert_markdown(textbox, text)
            else:
                textbox.insert("end", text)
            self._autosize_textbox(textbox)

        textbox.configure(state="disabled")
        self._scroll_to_bottom()
        return textbox

    def _configure_markdown_tags(self, textbox):
        mono = "SF Mono" if IS_MACOS else "Cascadia Mono" if IS_WINDOWS else "DejaVu Sans Mono"
        text_color = resolve_color(COLORS["text"])
        muted = resolve_color(COLORS["muted"])
        accent = resolve_color(COLORS["accent"])
        code_bg = resolve_color(COLORS["surface_alt"])
        textbox.tag_configure("bold", font=(UI_FONT_FAMILY, 11, "bold"))
        textbox.tag_configure("italic", font=(UI_FONT_FAMILY, 11, "italic"))
        textbox.tag_configure(
            "bold_italic", font=(UI_FONT_FAMILY, 11, "bold italic")
        )
        textbox.tag_configure("strike", overstrike=True, foreground=muted)
        textbox.tag_configure(
            "inline_code", font=(mono, 10), background=code_bg,
            foreground=text_color,
        )
        textbox.tag_configure(
            "code", font=(mono, 10), background=code_bg, foreground=text_color,
            lmargin1=12, lmargin2=12, rmargin=12, spacing1=5, spacing3=5,
        )
        textbox.tag_configure(
            "code_language", font=(mono, 8, "bold"), foreground=muted,
            background=code_bg, lmargin1=12, lmargin2=12, rmargin=12,
            spacing1=5,
        )
        textbox.tag_configure(
            "h1", font=(UI_FONT_FAMILY, 18, "bold"), spacing1=8, spacing3=5
        )
        textbox.tag_configure(
            "h2", font=(UI_FONT_FAMILY, 16, "bold"), spacing1=7, spacing3=4
        )
        textbox.tag_configure(
            "h3", font=(UI_FONT_FAMILY, 14, "bold"), spacing1=6, spacing3=3
        )
        textbox.tag_configure(
            "h4", font=(UI_FONT_FAMILY, 12, "bold"), spacing1=5, spacing3=2
        )
        textbox.tag_configure("list", lmargin1=16, lmargin2=30, spacing1=2)
        textbox.tag_configure(
            "list_marker", foreground=accent, font=(UI_FONT_FAMILY, 11, "bold")
        )
        textbox.tag_configure(
            "quote", foreground=muted, lmargin1=16, lmargin2=16,
            rmargin=8, spacing1=3, spacing3=3,
        )
        textbox.tag_configure(
            "rule", foreground=resolve_color(COLORS["border"]), spacing1=4, spacing3=4
        )
        textbox.tag_configure(
            "table", font=(mono, 9), background=code_bg,
            lmargin1=8, lmargin2=8, rmargin=8,
        )
        textbox.tag_configure("error", foreground=resolve_color(COLORS["danger"]))

    def _insert_markdown(self, textbox, text):
        """Render common GitHub-style Markdown into a read-only Tk text surface."""
        lines = text.splitlines()
        index = 0
        while index < len(lines):
            line = lines[index]
            fence = re.match(r"^\s*(```|''')\s*([\w.+-]*)\s*$", line)
            if fence:
                marker, language = fence.groups()
                code_lines = []
                index += 1
                while index < len(lines) and not re.match(
                    rf"^\s*{re.escape(marker)}\s*$", lines[index]
                ):
                    code_lines.append(lines[index])
                    index += 1
                if language:
                    textbox.insert("end", language.upper() + "\n", "code_language")
                textbox.insert("end", "\n".join(code_lines) + "\n", "code")
                index += 1
                continue

            heading = re.match(r"^(#{1,6})\s+(.+?)\s*#*\s*$", line)
            if heading:
                level = min(len(heading.group(1)), 4)
                self._insert_markdown_inline(
                    textbox, heading.group(2), (f"h{level}",)
                )
                textbox.insert("end", "\n", f"h{level}")
                index += 1
                continue

            if re.match(r"^\s{0,3}([-*_])(?:\s*\1){2,}\s*$", line):
                textbox.insert("end", "\u2500" * 42 + "\n", "rule")
                index += 1
                continue

            # Render Markdown tables as aligned, readable monospace rows.
            if (
                "|" in line and index + 1 < len(lines)
                and re.match(r"^\s*\|?\s*:?-{3,}", lines[index + 1])
            ):
                table_lines = [line]
                index += 2  # Skip the Markdown alignment separator.
                while index < len(lines) and "|" in lines[index] and lines[index].strip():
                    table_lines.append(lines[index])
                    index += 1
                textbox.insert("end", "\n".join(table_lines) + "\n", "table")
                continue

            list_item = re.match(r"^(\s*)([-+*]|\d+[.)])\s+(.+)$", line)
            if list_item:
                indent, marker, content = list_item.groups()
                task = re.match(r"^\[([ xX])\]\s+(.+)$", content)
                if task:
                    marker = "\u2611" if task.group(1).lower() == "x" else "\u2610"
                    content = task.group(2)
                elif not marker[0].isdigit():
                    marker = "\u2022"
                textbox.insert("end", f"{indent}{marker}  ", ("list", "list_marker"))
                self._insert_markdown_inline(textbox, content, ("list",))
                textbox.insert("end", "\n", "list")
                index += 1
                continue

            quote = re.match(r"^\s*>\s?(.*)$", line)
            if quote:
                textbox.insert("end", "\u2502  ", ("quote", "list_marker"))
                self._insert_markdown_inline(textbox, quote.group(1), ("quote",))
                textbox.insert("end", "\n", "quote")
                index += 1
                continue

            self._insert_markdown_inline(textbox, line)
            textbox.insert("end", "\n")
            index += 1

    def _insert_markdown_inline(self, textbox, text, base_tags=()):
        token_pattern = re.compile(
            r"(!?\[[^\]]+\]\(https?://[^\s)]+\)"
            r"|\*\*\*.+?\*\*\*|___.+?___"
            r"|\*\*.+?\*\*|__.+?__|~~.+?~~|`[^`]+`"
            r"|(?<!\*)\*[^*\n]+\*(?!\*)|(?<!_)_[^_\n]+_(?!_)"
            r"|https?://[^\s<>\]\)]+)"
        )
        position = 0
        for match in token_pattern.finditer(text):
            if match.start() > position:
                textbox.insert("end", text[position:match.start()], base_tags)
            token = match.group(0)
            tags = base_tags
            if token.startswith("![") or token.startswith("["):
                image_link = token.startswith("![")
                parsed = re.match(r"!?\[([^\]]+)\]\((https?://[^\s)]+)\)", token)
                if parsed:
                    label, url = parsed.groups()
                    self._insert_link(
                        textbox, f"Image: {label}" if image_link else label,
                        url, base_tags,
                    )
            elif token.startswith(("***", "___")):
                textbox.insert("end", token[3:-3], tags + ("bold_italic",))
            elif token.startswith(("**", "__")):
                textbox.insert("end", token[2:-2], tags + ("bold",))
            elif token.startswith("~~"):
                textbox.insert("end", token[2:-2], tags + ("strike",))
            elif token.startswith("`"):
                textbox.insert("end", token[1:-1], tags + ("inline_code",))
            elif token.startswith(("*", "_")):
                textbox.insert("end", token[1:-1], tags + ("italic",))
            elif token.startswith(("http://", "https://")):
                self._insert_link(textbox, token, token, base_tags)
            position = match.end()
        if position < len(text):
            textbox.insert("end", text[position:], base_tags)

    def _insert_link(self, textbox, label, url, base_tags=()):
        counter = getattr(self, "_link_counter", 0) + 1
        self._link_counter = counter
        tag_name = f"link_{counter}"
        textbox.tag_configure(
            tag_name, foreground=resolve_color(COLORS["accent"]), underline=True
        )
        textbox.tag_bind(
            tag_name, "<Button-1>", lambda _event, target=url: webbrowser.open(target)
        )
        textbox.tag_bind(
            tag_name, "<Enter>", lambda _event, tb=textbox: tb.configure(cursor="hand2")
        )
        textbox.tag_bind(
            tag_name, "<Leave>", lambda _event, tb=textbox: tb.configure(cursor="arrow")
        )
        textbox.insert("end", label, base_tags + (tag_name,))

    def _autosize_textbox(self, textbox):
        """Auto-resize a Text widget to fit its content."""
        textbox.update_idletasks()
        # Count display lines (accounts for word wrap)
        try:
            count = textbox.count("1.0", "end", "displaylines")
            if count and count[0] > 0:
                # Larger heading fonts and block spacing consume more pixels than
                # a normal Text line, so reserve a few extra line units for them.
                extra = 0
                for tag_name in ("h1", "h2", "h3", "h4", "code", "table"):
                    extra += len(textbox.tag_ranges(tag_name)) // 2
                textbox.configure(height=count[0] + extra)
            else:
                # Fallback: count newlines
                num_lines = int(textbox.index("end-1c").split(".")[0])
                textbox.configure(height=max(1, num_lines))
        except Exception:
            num_lines = int(textbox.index("end-1c").split(".")[0])
            textbox.configure(height=max(1, num_lines))

    def add_system_bubble(self, text):
        outer = ctk.CTkFrame(self.chat_frame, fg_color="transparent")
        outer.pack(fill="x", padx=12, pady=(8, 4))

        label = ctk.CTkLabel(
            outer, text=text, wraplength=500, justify="center",
            text_color=COLORS["muted"], font=ui_font(12),
        )
        label.pack(pady=4)

    def _scroll_to_bottom(self):
        if not self._closed:
            self.chat_frame.after(
                80, lambda: self.chat_frame._parent_canvas.yview_moveto(1.0)
            )

    def send_message(self):
        text = self.input_var.get().strip()
        if not text or self.streaming:
            return

        self.input_var.set("")
        self.add_user_bubble(text)
        self.messages.append({"role": "user", "content": text})
        self._save_session()

        self.streaming = True
        self.send_btn.configure(state="disabled")
        self._current_ai_label = self.add_ai_bubble("\u2026")

        threading.Thread(target=self._stream_response, daemon=True).start()

    def _stream_response(self):
        try:
            response = self.client.chat.completions.create(
                model=self.model_id,
                messages=self.messages,
                stream=True,
                temperature=0.7,
            )

            full_text = ""
            for chunk in response:
                if self._closed:
                    return
                if chunk.choices and chunk.choices[0].delta.content:
                    full_text += chunk.choices[0].delta.content
                    text_snapshot = full_text
                    self.win.after(0, lambda t=text_snapshot: self._update_ai_label(t))

            self.messages.append({"role": "assistant", "content": full_text})
            self._save_session()
            if not self._closed:
                self.win.after(0, self._stream_done)
        except Exception as e:
            error_msg = str(e)
            if not self._closed:
                self.win.after(0, lambda err=error_msg: self._stream_error(err))

    def _update_ai_label(self, text):
        tb = self._current_ai_label
        if tb and tb.winfo_exists():
            tb.configure(state="normal")
            tb.delete("1.0", "end")
            tb.insert("1.0", text)
            self._autosize_textbox(tb)
            tb.configure(state="disabled")
            self._scroll_to_bottom()

    def _stream_done(self):
        self.streaming = False
        self.send_btn.configure(state="normal")
        # Replace the streaming text with fully rendered Markdown.
        tb = self._current_ai_label
        if tb and tb.winfo_exists():
            tb.configure(state="normal")
            raw_text = tb.get("1.0", "end-1c")
            tb.delete("1.0", "end")
            self._insert_markdown(tb, raw_text)
            self._autosize_textbox(tb)
            tb.configure(state="disabled")
        self.input_entry.focus()

    def _stream_error(self, error):
        self.streaming = False
        self.send_btn.configure(state="normal")
        tb = self._current_ai_label
        if tb and tb.winfo_exists():
            tb.configure(state="normal")
            tb.delete("1.0", "end")
            tb.insert("1.0", f"Error: {error}", "error")
            self._autosize_textbox(tb)
            tb.configure(state="disabled")


class CompletionPanel:
    """Translation result summary rendered inside the main app window."""

    def __init__(self, app, info):
        self.app = app
        self.info = info
        self._closing = False
        self.overlay = ctk.CTkFrame(
            app.master, corner_radius=0, fg_color=("#E9EDF4", "#090B0E")
        )
        self.overlay.place(x=0, y=0, relwidth=1, relheight=1)
        self.overlay.lift()

        self.card = ctk.CTkFrame(
            self.overlay, width=600, corner_radius=12,
            fg_color=COLORS["surface"], border_width=1,
            border_color=COLORS["border"],
        )
        self.card.place(relx=0.5, rely=0.52, relheight=0.9, anchor="center")
        self.card.grid_propagate(False)
        self.card.grid_columnconfigure(0, weight=1)
        self.card.grid_rowconfigure(1, weight=1)
        self._build_ui()
        self._animate_open(0)

    def _build_ui(self):
        header = ctk.CTkFrame(self.card, height=78, fg_color="transparent")
        header.grid(row=0, column=0, sticky="ew", padx=24, pady=(18, 4))
        header.grid_columnconfigure(1, weight=1)
        ctk.CTkLabel(
            header, text="", image=load_icon("check-circle", 28),
            width=42, height=42, corner_radius=21,
            fg_color=("#E0F5EB", "#183C2F"),
        ).grid(row=0, column=0, rowspan=2, padx=(0, 13))
        ctk.CTkLabel(
            header, text="Translation complete", anchor="w",
            text_color=COLORS["text"], font=ui_font(19, "bold"),
        ).grid(row=0, column=1, sticky="sw")
        ctk.CTkLabel(
            header, text="Your translated documents are ready.", anchor="w",
            text_color=COLORS["muted"], font=ui_font(10),
        ).grid(row=1, column=1, sticky="nw")
        close_btn = ctk.CTkButton(
            header, text="", image=load_icon("x", 17), width=32, height=32,
            corner_radius=6, fg_color="transparent", hover_color=COLORS["border"],
            command=self.close,
        )
        close_btn.grid(row=0, column=2, rowspan=2, padx=(12, 0))
        Tooltip(close_btn, "Close")

        body = ctk.CTkScrollableFrame(
            self.card, fg_color="transparent", corner_radius=0,
            scrollbar_button_color=COLORS["border"],
        )
        body.grid(row=1, column=0, sticky="nsew", padx=24, pady=4)
        body.grid_columnconfigure(0, weight=1)

        source_languages = list(dict.fromkeys(self.info.get("source_languages", [])))
        source_summary = ", ".join(source_languages) if source_languages else AUTO_DETECT
        model_label = MODELS.get(self.info.get("model_id"), {}).get(
            "label", self.info.get("model_id", "Unknown")
        )
        rows = [
            ("Files", f"{self.info['completed_files']}/{self.info['total_files']}"),
            ("From", source_summary),
            ("To", self.info.get("target_language", "English")),
            ("Model", model_label),
            ("Paragraphs", f"{self.info['paragraphs']:,}"),
            ("Tokens", f"{self.info['tokens']:,}"),
            ("Cost", f"${self.info['cost']:.4f}"),
            ("Time", f"{self.info['elapsed']}s"),
            ("Output", soft_wrap_path(self.info.get("output_dir", ""))),
        ]
        stats = ctk.CTkFrame(
            body, fg_color=COLORS["surface_alt"], corner_radius=8
        )
        stats.grid(row=0, column=0, sticky="ew", pady=(0, 14))
        stats.grid_columnconfigure(1, weight=1)
        for index, (label, value) in enumerate(rows):
            ctk.CTkLabel(
                stats, text=label, width=88, anchor="w",
                text_color=COLORS["muted"], font=ui_font(10),
            ).grid(row=index, column=0, sticky="nw", padx=(14, 8), pady=(8, 0))
            ctk.CTkLabel(
                stats, text=value, anchor="w", wraplength=410,
                text_color=COLORS["text"], font=ui_font(10, "bold"),
            ).grid(row=index, column=1, sticky="nw", padx=(0, 14), pady=(8, 0))

        file_results = self.info.get("file_results", [])
        if file_results:
            ctk.CTkLabel(
                body, text="Translated files", anchor="w",
                text_color=COLORS["text"], font=ui_font(11, "bold"),
            ).grid(row=1, column=0, sticky="ew", pady=(0, 7))
            for row_index, result in enumerate(file_results, 2):
                output_path = result.get("output_path", "")
                row = ctk.CTkFrame(body, fg_color="transparent", corner_radius=7)
                row.grid(row=row_index, column=0, sticky="ew", pady=2)
                row.grid_columnconfigure(0, weight=1)
                ctk.CTkLabel(
                    row, text=os.path.basename(output_path), anchor="w",
                    text_color=COLORS["text"], font=ui_font(10, "bold"),
                ).grid(row=0, column=0, sticky="ew", padx=(8, 6), pady=7)
                ctk.CTkButton(
                    row, text="Ask AI", width=64, height=28, corner_radius=6,
                    fg_color="transparent", hover_color=COLORS["surface_alt"],
                    border_width=1, border_color=COLORS["border"],
                    text_color=COLORS["accent"], font=ui_font(9, "bold"),
                    command=lambda path=output_path: self._ask_ai(path),
                ).grid(row=0, column=1, padx=(4, 0), pady=4)
                ctk.CTkButton(
                    row, text="Open", width=56, height=28, corner_radius=6,
                    fg_color="transparent", hover_color=COLORS["surface_alt"],
                    text_color=COLORS["text"], font=ui_font(9),
                    command=lambda path=output_path: self.app.open_file(path),
                ).grid(row=0, column=2, padx=(2, 4), pady=4)

        footer = ctk.CTkFrame(self.card, height=66, fg_color="transparent")
        footer.grid(row=2, column=0, sticky="ew", padx=24, pady=(8, 18))
        footer.grid_columnconfigure(0, weight=1)
        ctk.CTkButton(
            footer, text="Show in folder", height=38, corner_radius=7,
            fg_color="transparent", hover_color=COLORS["surface_alt"],
            border_width=1, border_color=COLORS["border"],
            text_color=COLORS["text"], font=ui_font(10),
            command=lambda: self.app.open_in_folder(self.info.get("output_dir", "")),
        ).grid(row=0, column=0, sticky="w")
        ctk.CTkButton(
            footer, text="View history", width=104, height=38, corner_radius=7,
            fg_color="transparent", hover_color=COLORS["surface_alt"],
            text_color=COLORS["text"], font=ui_font(10, "bold"),
            command=self._view_history,
        ).grid(row=0, column=1, padx=8)
        ctk.CTkButton(
            footer, text="Done", width=88, height=38, corner_radius=7,
            fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
            font=ui_font(10, "bold"), command=self.close,
        ).grid(row=0, column=2)

    def _animate_open(self, step):
        if not self.overlay.winfo_exists():
            return
        progress = min(1.0, step / 8)
        eased = 1 - (1 - progress) ** 3
        self.card.place_configure(rely=0.5 + 0.02 * (1 - eased))
        if step < 8:
            self.overlay.after(14, lambda: self._animate_open(step + 1))

    def close(self, after_close=None):
        if self._closing:
            return
        self._closing = True
        self._after_close = after_close
        self._animate_close(0)

    def _animate_close(self, step):
        if not self.overlay.winfo_exists():
            return
        progress = min(1.0, step / 7)
        self.card.place_configure(rely=0.5 + 0.025 * progress ** 2)
        if step < 7:
            self.overlay.after(14, lambda: self._animate_close(step + 1))
            return
        self.overlay.destroy()
        self.app.completion_panel = None
        if self._after_close:
            self._after_close()

    def _ask_ai(self, file_path):
        if not self.app.api_key_var.get().strip():
            self.app.open_settings(focus_api=True)
            return
        self.close(
            lambda: self.app.open_chat(file_path, self.info.get("model_id"))
        )

    def _view_history(self):
        self.close(self.app.open_history)


class HistoryPanel:
    """Durable translation history browser embedded in the main window."""

    def __init__(self, app):
        self.app = app
        self.selected_id = None
        self.overlay = ctk.CTkFrame(
            app.master, corner_radius=0, fg_color=COLORS["background"]
        )
        self.overlay.place(x=0, y=0, relwidth=1, relheight=1)
        self.overlay.lift()
        self._build_ui()
        self._render_history()

    def _build_ui(self):
        header = ctk.CTkFrame(
            self.overlay, height=62, corner_radius=0, fg_color=COLORS["surface"]
        )
        header.grid(row=0, column=0, sticky="ew")
        header.grid_propagate(False)
        header.grid_columnconfigure(1, weight=1)
        ctk.CTkButton(
            header, text="Back", width=64, height=32, corner_radius=6,
            fg_color="transparent", hover_color=COLORS["surface_alt"],
            text_color=COLORS["text"], font=ui_font(10, "bold"),
            command=self.close,
        ).grid(row=0, column=0, padx=(16, 8), pady=15)
        ctk.CTkLabel(
            header, text="Translation history", anchor="w",
            text_color=COLORS["text"], font=ui_font(17, "bold"),
        ).grid(row=0, column=1, sticky="w", pady=15)
        self.clear_btn = ctk.CTkButton(
            header, text="Clear history", width=92, height=32, corner_radius=6,
            fg_color="transparent", hover_color=("#FBE9EC", "#44242B"),
            text_color=COLORS["danger"], font=ui_font(9),
            command=self._clear_history,
        )
        self.clear_btn.grid(row=0, column=2, padx=16, pady=15)

        self.overlay.grid_columnconfigure(0, weight=1)
        self.overlay.grid_rowconfigure(1, weight=1)
        content = ctk.CTkFrame(self.overlay, fg_color="transparent")
        content.grid(row=1, column=0, sticky="nsew", padx=20, pady=20)
        content.grid_columnconfigure(1, weight=1)
        content.grid_rowconfigure(0, weight=1)

        left = ctk.CTkFrame(
            content, width=255, fg_color=COLORS["surface"], corner_radius=9,
            border_width=1, border_color=COLORS["border"],
        )
        left.grid(row=0, column=0, sticky="nsw", padx=(0, 12))
        left.grid_propagate(False)
        left.grid_rowconfigure(1, weight=1)
        left.grid_columnconfigure(0, weight=1)
        self.count_label = ctk.CTkLabel(
            left, text="", anchor="w", text_color=COLORS["muted"],
            font=ui_font(10, "bold"),
        )
        self.count_label.grid(row=0, column=0, sticky="ew", padx=14, pady=(13, 7))
        self.history_list = ctk.CTkScrollableFrame(
            left, fg_color="transparent", corner_radius=0,
            scrollbar_button_color=COLORS["border"],
        )
        self.history_list.grid(row=1, column=0, sticky="nsew", padx=6, pady=(0, 8))

        self.detail_host = ctk.CTkFrame(
            content, fg_color=COLORS["surface"], corner_radius=9,
            border_width=1, border_color=COLORS["border"],
        )
        self.detail_host.grid(row=0, column=1, sticky="nsew")
        self.detail_host.grid_columnconfigure(0, weight=1)
        self.detail_host.grid_rowconfigure(0, weight=1)

    def _render_history(self):
        for child in self.history_list.winfo_children():
            child.destroy()
        history = self.app.history
        self.count_label.configure(
            text=f"{len(history)} translation{'s' if len(history) != 1 else ''}"
        )
        self.clear_btn.configure(state="normal" if history else "disabled")
        if not history:
            self._show_empty()
            return
        for record in history:
            record_id = record.get("id", "")
            selected = record_id == self.selected_id
            count = record.get("completed_files", 0)
            target = record.get("target_language", "")
            text = (
                f"{format_history_time(record.get('created_at'))}\n"
                f"{count} file{'s' if count != 1 else ''}  \u00b7  {target}"
            )
            button = ctk.CTkButton(
                self.history_list, text=text, anchor="w", height=58,
                corner_radius=7,
                fg_color=("#EEF0FF", "#292E52") if selected else "transparent",
                hover_color=COLORS["surface_alt"],
                text_color=COLORS["accent"] if selected else COLORS["text"],
                font=ui_font(9, "bold" if selected else "normal"),
                command=lambda item=record: self.show_record(item),
            )
            button.pack(fill="x", padx=2, pady=2)
        if self.selected_id is None:
            self.show_record(history[0])

    def _show_empty(self):
        for child in self.detail_host.winfo_children():
            child.destroy()
        ctk.CTkLabel(
            self.detail_host, text="No translation history yet",
            text_color=COLORS["text"], font=ui_font(14, "bold"),
        ).place(relx=0.5, rely=0.46, anchor="center")
        ctk.CTkLabel(
            self.detail_host,
            text="Completed translations will appear here with their files and details.",
            text_color=COLORS["muted"], font=ui_font(10), wraplength=340,
        ).place(relx=0.5, rely=0.53, anchor="center")

    def show_record(self, record):
        self.selected_id = record.get("id")
        for child in self.detail_host.winfo_children():
            child.destroy()
        detail = ctk.CTkScrollableFrame(
            self.detail_host, fg_color="transparent", corner_radius=0,
            scrollbar_button_color=COLORS["border"],
        )
        detail.grid(row=0, column=0, sticky="nsew", padx=18, pady=16)
        detail.grid_columnconfigure(0, weight=1)

        ctk.CTkLabel(
            detail, text=format_history_time(record.get("created_at")), anchor="w",
            text_color=COLORS["text"], font=ui_font(15, "bold"),
        ).grid(row=0, column=0, sticky="ew")
        model_label = MODELS.get(record.get("model_id"), {}).get(
            "label", record.get("model_id", "Unknown")
        )
        source = ", ".join(record.get("source_languages", [])) or AUTO_DETECT
        stats = [
            ("Files", f"{record.get('completed_files', 0)}/{record.get('total_files', 0)}"),
            ("From", source),
            ("To", record.get("target_language", "")),
            ("Model", model_label),
            ("Paragraphs", f"{record.get('paragraphs', 0):,}"),
            ("Tokens", f"{record.get('tokens', 0):,}"),
            ("Cost", f"${record.get('cost', 0):.4f}"),
            ("Time", f"{record.get('elapsed', 0)}s"),
            ("Output folder", soft_wrap_path(record.get("output_dir", ""))),
        ]
        stat_frame = ctk.CTkFrame(
            detail, fg_color=COLORS["surface_alt"], corner_radius=8
        )
        stat_frame.grid(row=1, column=0, sticky="ew", pady=(12, 16))
        stat_frame.grid_columnconfigure(1, weight=1)
        for index, (label, value) in enumerate(stats):
            ctk.CTkLabel(
                stat_frame, text=label, width=90, anchor="w",
                text_color=COLORS["muted"], font=ui_font(9),
            ).grid(row=index, column=0, sticky="nw", padx=(12, 8), pady=(7, 0))
            ctk.CTkLabel(
                stat_frame, text=value, anchor="w", wraplength=330,
                text_color=COLORS["text"], font=ui_font(9, "bold"),
            ).grid(row=index, column=1, sticky="nw", padx=(0, 12), pady=(7, 0))

        ctk.CTkLabel(
            detail, text="Files", anchor="w", text_color=COLORS["text"],
            font=ui_font(11, "bold"),
        ).grid(row=2, column=0, sticky="ew", pady=(0, 6))
        file_results = record.get("file_results", [])
        for index, result in enumerate(file_results, 3):
            self._add_file_detail(detail, index, result, record)

        ctk.CTkButton(
            detail, text="Remove this record", width=116, height=30,
            corner_radius=6, fg_color="transparent",
            hover_color=("#FBE9EC", "#44242B"), text_color=COLORS["danger"],
            font=ui_font(9), command=lambda: self._remove_record(record),
        ).grid(row=3 + len(file_results), column=0, sticky="w", pady=(14, 4))
        self._render_history_selection_only()

    def _render_history_selection_only(self):
        for child, record in zip(self.history_list.winfo_children(), self.app.history):
            selected = record.get("id") == self.selected_id
            child.configure(
                fg_color=("#EEF0FF", "#292E52") if selected else "transparent",
                text_color=COLORS["accent"] if selected else COLORS["text"],
            )

    def _add_file_detail(self, parent, row_index, result, record):
        input_path = result.get("input_path", "")
        output_path = result.get("output_path", "")
        exists = Path(output_path).exists()
        card = ctk.CTkFrame(
            parent, fg_color="transparent", corner_radius=7,
            border_width=1, border_color=COLORS["border"],
        )
        card.grid(row=row_index, column=0, sticky="ew", pady=4)
        card.grid_columnconfigure(0, weight=1)
        ctk.CTkLabel(
            card, text=os.path.basename(output_path) or "Missing output", anchor="w",
            text_color=COLORS["text"], font=ui_font(10, "bold"),
        ).grid(row=0, column=0, columnspan=3, sticky="ew", padx=10, pady=(8, 1))
        ctk.CTkLabel(
            card,
            text=(
                f"From: {soft_wrap_path(input_path)}\n"
                f"To: {soft_wrap_path(output_path)}"
            ), anchor="w",
            justify="left", wraplength=390, text_color=COLORS["muted"],
            font=ui_font(8),
        ).grid(row=1, column=0, columnspan=3, sticky="ew", padx=10, pady=(0, 6))
        ask = ctk.CTkButton(
            card, text="Ask AI", width=66, height=28, corner_radius=6,
            fg_color=COLORS["accent"] if exists else COLORS["surface_alt"],
            hover_color=COLORS["accent_hover"], font=ui_font(9, "bold"),
            state="normal" if exists else "disabled",
            command=lambda: self.app.open_chat(output_path, record.get("model_id")),
        )
        ask.grid(row=2, column=0, sticky="w", padx=(10, 4), pady=(0, 9))
        ctk.CTkButton(
            card, text="Open", width=54, height=28, corner_radius=6,
            fg_color="transparent", hover_color=COLORS["surface_alt"],
            text_color=COLORS["text"], font=ui_font(9),
            state="normal" if exists else "disabled",
            command=lambda: self.app.open_file(output_path),
        ).grid(row=2, column=1, sticky="w", padx=2, pady=(0, 9))
        ctk.CTkButton(
            card, text="Reveal", width=56, height=28, corner_radius=6,
            fg_color="transparent", hover_color=COLORS["surface_alt"],
            text_color=COLORS["text"], font=ui_font(9),
            state="normal" if exists else "disabled",
            command=lambda: self.app.reveal_file(output_path),
        ).grid(row=2, column=2, sticky="w", padx=2, pady=(0, 9))

    def _remove_record(self, record):
        self.app.history = [
            item for item in self.app.history if item.get("id") != record.get("id")
        ]
        save_translation_history(self.app.history)
        self.selected_id = None
        self._render_history()

    def _clear_history(self):
        if not messagebox.askyesno(
            "Clear translation history",
            "Remove all saved translation records? Your translated files and chats will not be deleted.",
            parent=self.app.master,
        ):
            return
        self.app.history = []
        save_translation_history([])
        self.selected_id = None
        self._render_history()

    def close(self):
        self.app.history_panel = None
        if self.overlay.winfo_exists():
            self.overlay.destroy()


class PPTTranslatorApp:
    def __init__(self, master):
        self.master = master
        self.config = load_config()
        appearance = self.config.get("appearance", "System")
        if appearance not in {"System", "Light", "Dark"}:
            appearance = "System"

        ctk.set_appearance_mode(appearance)
        ctk.set_default_color_theme("blue")

        master.title(APP_DISPLAY_NAME)
        master.minsize(720, 560)
        master.configure(fg_color=COLORS["background"])
        set_window_icon(master)
        center_window(master, 900, 680)

        self.file_queue = []
        self.file_widgets = []
        self.file_widget_map = {}
        self.file_states = {}
        self.cancel_event = threading.Event()
        self.translating = False
        self._swap_animating = False
        self._theme_switching = False
        self.detected_source_languages = []
        self.settings_window = None
        self.completion_panel = None
        self.history_panel = None
        self.chat_panel = None
        self.history = load_translation_history()
        self._notice_job = None
        self.appearance_var = ctk.StringVar(value=appearance)

        self.api_key_var = ctk.StringVar(value=self.config.get("api_key", ""))
        self.model_ids = list(MODELS.keys())
        model_labels = [MODELS[mid]["label"] for mid in self.model_ids]
        saved_model = self.config.get("model", MODELS[self.model_ids[1]]["label"])
        if saved_model not in model_labels:
            saved_model = MODELS[self.model_ids[1]]["label"]
        self.model_var = ctk.StringVar(value=saved_model)

        saved_source = self.config.get("source_language", AUTO_DETECT)
        if saved_source not in (AUTO_DETECT, *LANGUAGES):
            saved_source = AUTO_DETECT
        self.source_language_var = ctk.StringVar(value=saved_source)

        saved_target = self.config.get("target_language", "English")
        if saved_target not in LANGUAGES:
            saved_target = "English"
        self.target_language_var = ctk.StringVar(value=saved_target)

        self.setup_ui()
        self.register_shortcuts()
        self.register_drag_and_drop()
        self.master.protocol("WM_DELETE_WINDOW", self.on_close)

    def setup_ui(self):
        self.master.grid_columnconfigure(0, weight=1)
        self.master.grid_rowconfigure(2, weight=1)

        # Compact application header.
        header = ctk.CTkFrame(
            self.master, height=62, corner_radius=0, fg_color=COLORS["surface"]
        )
        header.grid(row=0, column=0, sticky="ew")
        header.grid_propagate(False)
        header.grid_columnconfigure(1, weight=1)

        ctk.CTkLabel(
            header, text="", image=load_brand_image(30), width=34, height=34,
        ).grid(row=0, column=0, padx=(22, 10), pady=14)
        ctk.CTkLabel(
            header, text=APP_DISPLAY_NAME, anchor="w", text_color=COLORS["text"],
            font=ui_font(17, "bold"),
        ).grid(row=0, column=1, sticky="w", pady=14)
        self.history_btn = ctk.CTkButton(
            header, text="History", width=68, height=34, corner_radius=7,
            fg_color="transparent", hover_color=COLORS["border"],
            text_color=COLORS["text"], font=ui_font(10, "bold"),
            command=self.open_history,
        )
        self.history_btn.grid(row=0, column=2, padx=(10, 2), pady=14)
        self.settings_btn = ctk.CTkButton(
            header, text="", image=load_icon("settings", 19), width=36, height=36,
            corner_radius=7, fg_color="transparent", hover_color=COLORS["border"],
            command=self.open_settings,
        )
        self.settings_btn.grid(row=0, column=3, padx=(4, 20), pady=13)
        Tooltip(self.settings_btn, "Settings")
        ctk.CTkFrame(
            self.master, height=1, corner_radius=0, fg_color=COLORS["border"]
        ).grid(row=0, column=0, sticky="sew")

        # Source and target controls stay immediately above the work area.
        toolbar = ctk.CTkFrame(
            self.master, height=70, corner_radius=0, fg_color="transparent"
        )
        toolbar.grid(row=1, column=0, sticky="ew", padx=24)
        toolbar.grid_propagate(False)
        toolbar.grid_columnconfigure(0, weight=1)
        toolbar.grid_columnconfigure(7, weight=1)

        ctk.CTkLabel(
            toolbar, text="From", text_color=COLORS["muted"],
            font=ui_font(10, "bold"),
        ).grid(row=0, column=1, sticky="w", padx=(0, 7), pady=16)
        self.source_language_menu = ModernComboBox(
            toolbar, variable=self.source_language_var,
            values=[AUTO_DETECT, *LANGUAGES], width=250, height=40,
            max_visible=7, searchable=True,
            command=self.on_source_language_change,
        )
        self.source_language_menu.grid(row=0, column=2, sticky="ew", pady=15)

        self.swap_language_btn = ctk.CTkButton(
            toolbar, text="", image=load_icon("arrow-left-right", 18),
            width=38, height=38, corner_radius=7,
            fg_color=COLORS["surface"], hover_color=COLORS["border"],
            border_width=1, border_color=COLORS["border"],
            text_color_disabled=COLORS["muted"], command=self.swap_languages,
        )
        self.swap_language_btn.grid(row=0, column=3, padx=12, pady=16)
        Tooltip(self.swap_language_btn, "Swap source and target languages")

        ctk.CTkLabel(
            toolbar, text="To", text_color=COLORS["muted"],
            font=ui_font(10, "bold"),
        ).grid(row=0, column=4, sticky="w", padx=(0, 7), pady=16)
        self.target_language_menu = ModernComboBox(
            toolbar, variable=self.target_language_var, values=list(LANGUAGES),
            width=230, height=40, max_visible=7, searchable=True,
            command=self.on_target_language_change,
        )
        self.target_language_menu.grid(row=0, column=5, sticky="ew", pady=15)

        self.update_swap_state()

        # The document workspace is the dominant surface in the window.
        self.workspace_frame = ctk.CTkFrame(
            self.master, fg_color=COLORS["surface"], corner_radius=9,
            border_width=1, border_color=COLORS["border"],
        )
        self.workspace_frame.grid(
            row=2, column=0, sticky="nsew", padx=24, pady=(0, 18)
        )
        self.workspace_frame.grid_columnconfigure(0, weight=1)
        self.workspace_frame.grid_rowconfigure(1, weight=1)

        self.notice_label = ctk.CTkLabel(
            self.workspace_frame, text="", height=30, corner_radius=6,
            fg_color=COLORS["surface_alt"], text_color=COLORS["muted"],
            font=ui_font(10, "bold"),
        )
        self.notice_label.grid(row=0, column=0, sticky="ew", padx=14, pady=(12, 0))
        self.notice_label.grid_remove()

        self.empty_state = ctk.CTkFrame(
            self.workspace_frame, fg_color="transparent", corner_radius=0,
        )
        self.empty_state.grid(row=1, column=0, sticky="nsew", padx=24, pady=20)
        self.empty_state.grid_columnconfigure(0, weight=1)
        self.empty_state.grid_rowconfigure(0, weight=1)
        empty_content = ctk.CTkFrame(self.empty_state, fg_color="transparent")
        empty_content.grid(row=0, column=0)
        ctk.CTkLabel(
            empty_content, text="", image=load_icon("upload", 30),
        ).pack(pady=(0, 12))
        self.drop_title = ctk.CTkLabel(
            empty_content, text="Drop presentations here", text_color=COLORS["text"],
            font=ui_font(14, "bold"),
        )
        self.drop_title.pack(pady=(0, 5))
        ctk.CTkLabel(
            empty_content, text="PPTX  \u00b7  PDF", text_color=COLORS["muted"],
            font=ui_font(10),
        ).pack(pady=(0, 16))
        self.choose_files_btn = ctk.CTkButton(
            empty_content, text="Choose files", width=112, height=36, corner_radius=7,
            fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
            font=ui_font(11, "bold"), command=self.add_files,
        )
        self.choose_files_btn.pack()

        self.file_view = ctk.CTkFrame(
            self.workspace_frame, fg_color="transparent", corner_radius=0,
        )
        self.file_view.grid(row=1, column=0, sticky="nsew", padx=14, pady=(10, 8))
        self.file_view.grid_columnconfigure(0, weight=1)
        self.file_view.grid_rowconfigure(0, weight=1)
        self.files_frame = ctk.CTkScrollableFrame(
            self.file_view, fg_color="transparent", corner_radius=0,
            scrollbar_button_color=COLORS["border"],
            scrollbar_button_hover_color=COLORS["muted"],
        )
        self.files_frame.grid(row=0, column=0, sticky="nsew")
        self.files_frame.grid_columnconfigure(0, weight=1)
        self.add_files_btn = ctk.CTkButton(
            self.file_view, text="Add file", image=load_icon("plus", 15),
            compound="left", width=94, height=32, corner_radius=6,
            fg_color="transparent", hover_color=COLORS["surface_alt"],
            border_width=1, border_color=COLORS["border"],
            text_color=COLORS["text"], font=ui_font(10, "bold"),
            command=self.add_files,
        )
        self.add_files_btn.grid(row=1, column=0, sticky="w", padx=8, pady=(7, 3))
        self.file_view.grid_remove()

        # Global progress is a compact strip; file rows carry the detailed state.
        self.progress_frame = ctk.CTkFrame(
            self.workspace_frame, fg_color=COLORS["surface_alt"], corner_radius=7,
        )
        self.progress_frame.grid(row=2, column=0, sticky="ew", padx=14, pady=(0, 12))
        self.progress_frame.grid_columnconfigure(0, weight=1)

        self.progress_detail = ctk.CTkLabel(
            self.progress_frame, text="Preparing translation", anchor="w",
            text_color=COLORS["text"], font=ui_font(10, "bold"),
        )
        self.progress_detail.grid(row=0, column=0, sticky="ew", padx=(14, 8), pady=(10, 5))
        self.progress_percent = ctk.CTkLabel(
            self.progress_frame, text="0%", text_color=COLORS["accent"],
            font=ui_font(11, "bold"),
        )
        self.progress_percent.grid(row=0, column=1, padx=(8, 14), pady=(10, 5))

        self.progress_bar = ctk.CTkProgressBar(
            self.progress_frame, height=5, corner_radius=2,
            fg_color=COLORS["border"], progress_color=COLORS["accent"],
        )
        self.progress_bar.grid(row=1, column=0, columnspan=2, sticky="ew", padx=14)
        self.progress_bar.set(0)

        self.progress_stats = ctk.CTkLabel(
            self.progress_frame, text="", anchor="w", text_color=COLORS["muted"],
            font=ui_font(9),
        )
        self.progress_stats.grid(row=2, column=0, sticky="ew", padx=14, pady=(5, 10))
        self.cancel_btn = ctk.CTkButton(
            self.progress_frame, text="Cancel", width=68, height=28, corner_radius=6,
            fg_color="transparent", hover_color=("#FBE9EC", "#44242B"),
            text_color=COLORS["danger"], font=ui_font(11),
            command=self.cancel_translation,
        )
        self.cancel_btn.grid(row=2, column=1, padx=10, pady=(3, 8))
        self.progress_frame.grid_remove()

        # Sticky, minimal action bar.
        action_bar = ctk.CTkFrame(
            self.master, height=68, fg_color=COLORS["surface"], corner_radius=0,
        )
        action_bar.grid(row=3, column=0, sticky="ew")
        action_bar.grid_propagate(False)
        action_bar.grid_columnconfigure(0, weight=1)
        ctk.CTkFrame(
            action_bar, height=1, corner_radius=0, fg_color=COLORS["border"]
        ).grid(row=0, column=0, columnspan=2, sticky="new")

        self.status_label = ctk.CTkLabel(
            action_bar, text="", anchor="w",
            text_color=COLORS["muted"], font=ui_font(11),
        )
        self.status_label.grid(row=1, column=0, sticky="ew", padx=24, pady=(14, 12))

        self.translate_btn = ctk.CTkButton(
            action_bar, text="Translate", width=144, height=40, corner_radius=7,
            fg_color=COLORS["surface_alt"], hover_color=COLORS["surface_alt"],
            text_color="white", text_color_disabled=COLORS["muted"],
            font=ui_font(12, "bold"), command=self.start_translation, state="disabled",
        )
        self.translate_btn.grid(row=1, column=1, padx=24, pady=(13, 11))

    def register_shortcuts(self):
        modifier = "Command" if IS_MACOS else "Control"
        self.master.bind(f"<{modifier}-o>", lambda _event: self.add_files())
        self.master.bind(f"<{modifier}-Return>", lambda _event: self.start_translation())
        self.master.bind(f"<{modifier}-comma>", lambda _event: self.open_settings())
        if IS_MACOS:
            self.master.bind("<Command-q>", lambda _event: self.on_close())

    def register_drag_and_drop(self):
        if DND_FILES is None or not hasattr(self.master, "drop_target_register"):
            return
        try:
            self.master.drop_target_register(DND_FILES)
            self.master.dnd_bind("<<DropEnter>>", self._on_drop_enter)
            self.master.dnd_bind("<<DropLeave>>", self._on_drop_leave)
            self.master.dnd_bind("<<Drop>>", self._on_drop)
        except tk.TclError:
            # The app remains fully usable through the native file picker.
            pass

    def _on_drop_enter(self, _event):
        self.workspace_frame.configure(border_color=COLORS["accent"], border_width=2)
        if not self.file_queue:
            self.drop_title.configure(text="Drop to add files")
        else:
            self.notice_label.configure(
                text="Drop to add files", text_color=COLORS["accent"]
            )
            self.notice_label.grid()
        return "copy"

    def _on_drop_leave(self, _event):
        self._reset_drop_feedback()

    def _reset_drop_feedback(self):
        self.workspace_frame.configure(border_color=COLORS["border"], border_width=1)
        self.drop_title.configure(text="Drop presentations here")
        if self.notice_label.cget("text") == "Drop to add files":
            self.notice_label.grid_remove()

    def _on_drop(self, event):
        self._reset_drop_feedback()
        try:
            paths = list(self.master.tk.splitlist(event.data))
        except (tk.TclError, AttributeError):
            paths = []
        self.add_file_paths(paths)
        return "copy"

    def open_settings(self, focus_api=False):
        if self.translating:
            return
        if self.settings_window and self.settings_window.is_open:
            self.settings_window.backdrop.lift()
            if focus_api:
                self.settings_window.api_key_entry.focus_set()
            return
        self.settings_window = SettingsPanel(self, focus_api=focus_api)

    def open_history(self):
        if self.translating or (self.settings_window and self.settings_window.is_open):
            return
        if self.history_panel and self.history_panel.overlay.winfo_exists():
            self.history_panel.overlay.lift()
            return
        self.history_panel = HistoryPanel(self)

    def open_chat(self, file_path, model_id=None):
        path = str(Path(file_path).resolve())
        if not Path(path).is_file():
            messagebox.showerror(
                "File not found",
                "The selected translated file is no longer at its saved location.",
                parent=self.master,
            )
            return
        api_key = self.api_key_var.get().strip()
        if not api_key:
            self.open_settings(focus_api=True)
            return
        if self.chat_panel and self.chat_panel.overlay.winfo_exists():
            self.chat_panel.close()
        self.chat_panel = ChatPanel(
            self, path, model_id or self.get_selected_model_id(), api_key
        )

    def add_history_record(self, info):
        record = {
            "id": hashlib.sha256(
                f"{time.time_ns()}:{info.get('output_dir', '')}".encode("utf-8")
            ).hexdigest()[:16],
            "created_at": datetime.now(timezone.utc).isoformat(),
            "completed_files": int(info.get("completed_files", 0)),
            "total_files": int(info.get("total_files", 0)),
            "paragraphs": int(info.get("paragraphs", 0)),
            "tokens": int(info.get("tokens", 0)),
            "cost": float(info.get("cost", 0)),
            "elapsed": int(info.get("elapsed", 0)),
            "output_dir": info.get("output_dir", ""),
            "input_paths": list(info.get("input_paths", [])),
            "output_paths": list(info.get("output_paths", [])),
            "file_results": list(info.get("file_results", [])),
            "model_id": info.get("model_id", self.get_selected_model_id()),
            "source_languages": list(info.get("source_languages", [])),
            "target_language": info.get("target_language", "English"),
        }
        self.history.insert(0, record)
        self.history = self.history[:200]
        save_translation_history(self.history)
        info["history_id"] = record["id"]
        info["created_at"] = record["created_at"]
        return record

    def show_inline_notice(self, message, danger=False, duration=4000):
        if self._notice_job:
            self.master.after_cancel(self._notice_job)
            self._notice_job = None
        self.notice_label.configure(
            text=message,
            text_color=COLORS["danger"] if danger else COLORS["muted"],
        )
        self.notice_label.grid()

        def hide():
            self.notice_label.grid_remove()
            self._notice_job = None

        self._notice_job = self.master.after(duration, hide)

    def change_appearance_mode(self, appearance):
        if appearance not in {"System", "Light", "Dark"} or self._theme_switching:
            return
        self._theme_switching = True
        try:
            ctk.set_appearance_mode(appearance)
            self.config["appearance"] = appearance
            save_config(self.config)
        finally:
            self.master.after_idle(
                lambda: setattr(self, "_theme_switching", False)
            )

    def on_model_change(self, model_label):
        self.config["model"] = model_label
        save_config(self.config)

    def on_source_language_change(self, language):
        self.config["source_language"] = language
        self.detected_source_languages = []
        self.update_swap_state()
        save_config(self.config)

    def on_target_language_change(self, language):
        self.config["target_language"] = language
        save_config(self.config)

    def update_swap_state(self):
        can_swap = (
            not self.translating
            and not self._swap_animating
            and (
                self.source_language_var.get() != AUTO_DETECT
                or len(self.detected_source_languages) == 1
            )
        )
        self.swap_language_btn.configure(state="normal" if can_swap else "disabled")

    def swap_languages(self):
        if self._swap_animating:
            return
        source = self.source_language_var.get()
        target = self.target_language_var.get()
        if source == AUTO_DETECT:
            if len(self.detected_source_languages) != 1:
                return
            source = self.detected_source_languages[0]

        self._swap_animating = True
        self.swap_language_btn.configure(state="disabled")
        self._animate_language_swap(target, source, 0)

    def _animate_language_swap(self, new_source, new_target, step):
        midpoint = 5
        final_step = 11
        if step < midpoint:
            progress = step / midpoint
            self.source_language_menu.animate_content(
                offset=int(9 * progress), tone="muted"
            )
            self.target_language_menu.animate_content(
                offset=-int(7 * progress), tone="muted"
            )
            self.swap_language_btn.configure(
                fg_color=COLORS["surface_alt"],
                border_color=COLORS["accent"],
            )
        elif step == midpoint:
            self.source_language_var.set(new_source)
            self.target_language_var.set(new_target)
            self.detected_source_languages = []
            self.config["source_language"] = new_source
            self.config["target_language"] = new_target
            save_config(self.config)
            self.source_language_menu.animate_content(offset=-9, tone="accent")
            self.target_language_menu.animate_content(offset=7, tone="accent")
            self.swap_language_btn.configure(fg_color=COLORS["accent"])
        else:
            progress = (step - midpoint) / (final_step - midpoint)
            self.source_language_menu.animate_content(
                offset=int(-9 * (1 - progress)),
                tone="accent" if progress < 0.7 else "normal",
            )
            self.target_language_menu.animate_content(
                offset=int(7 * (1 - progress)),
                tone="accent" if progress < 0.7 else "normal",
            )
            if progress > 0.45:
                self.swap_language_btn.configure(fg_color=COLORS["surface"])

        if step < final_step:
            self.master.after(
                18,
                lambda: self._animate_language_swap(
                    new_source, new_target, step + 1
                ),
            )
        else:
            self.source_language_menu.animate_content(tone="normal")
            self.target_language_menu.animate_content(tone="normal")
            self.swap_language_btn.configure(
                fg_color=COLORS["surface"], border_color=COLORS["border"]
            )
            self._swap_animating = False
            self.update_swap_state()

    def show_detected_languages(self, detected_languages):
        unique_languages = list(dict.fromkeys(detected_languages))
        self.detected_source_languages = unique_languages
        self.update_swap_state()

    def set_translate_available(self, available):
        self.translate_btn.configure(
            state="normal" if available else "disabled",
            fg_color=COLORS["accent"] if available else COLORS["surface_alt"],
            hover_color=COLORS["accent_hover"] if available else COLORS["surface_alt"],
        )

    def set_controls_busy(self, busy):
        state = "disabled" if busy else "normal"
        for widget in (
            self.source_language_menu,
            self.target_language_menu,
            self.add_files_btn,
            self.choose_files_btn,
            self.settings_btn,
            self.history_btn,
        ):
            widget.configure(state=state)
        for entry in self.file_widget_map.values():
            entry.more_btn.configure(state=state)
        if busy:
            self.swap_language_btn.configure(state="disabled")
        else:
            self.update_swap_state()

    def on_close(self):
        if self.translating:
            should_close = messagebox.askyesno(
                "Translation in progress",
                "A translation is still running. Cancel it and close the app?",
                parent=self.master,
            )
            if not should_close:
                return
            self.cancel_event.set()
        if self.chat_panel:
            self.chat_panel._save_session()
        self.master.destroy()

    def save_api_key(self):
        self.config["api_key"] = self.api_key_var.get().strip()
        save_config(self.config)

    def get_selected_model_id(self):
        label = self.model_var.get()
        for mid, info in MODELS.items():
            if info["label"] == label:
                return mid
        return self.model_ids[1]

    def add_files(self):
        if self.settings_window and self.settings_window.is_open:
            return
        files = filedialog.askopenfilenames(
            parent=self.master,
            title="Choose presentations",
            initialdir=self.config.get("last_input_dir", str(Path.home())),
            filetypes=[
                ("Supported documents", ("*.pptx", "*.pdf")),
                ("PowerPoint", "*.pptx"),
                ("PDF", "*.pdf"),
            ],
        )
        if files:
            self.add_file_paths(files)

    def add_file_paths(self, paths):
        accepted = []
        unsupported = []
        for raw_path in paths:
            path = str(Path(raw_path).expanduser().resolve())
            if not os.path.isfile(path) or Path(path).suffix.lower() not in {".pptx", ".pdf"}:
                unsupported.append(path)
                continue
            if path not in self.file_queue:
                self.file_queue.append(path)
                self.file_states[path] = ("Ready", None, None)
                accepted.append(path)

        if accepted:
            self.config["last_input_dir"] = str(Path(accepted[0]).parent)
            save_config(self.config)
            self.refresh_file_list()
        if unsupported:
            self.show_inline_notice(
                "Unsupported file type \u2014 PPTX and PDF only", danger=True
            )
        elif paths and not accepted:
            self.show_inline_notice("Those files are already in the list")

    def remove_file(self, filepath):
        if self.translating:
            return
        if filepath in self.file_queue:
            self.file_queue.remove(filepath)
            self.file_states.pop(filepath, None)
            self.refresh_file_list()

    def refresh_file_list(self):
        for w in self.file_widgets:
            w.destroy()
        self.file_widgets.clear()
        self.file_widget_map.clear()

        if not self.file_queue:
            self.file_view.grid_remove()
            self.empty_state.grid()
            self.update_action_bar()
            return

        self.empty_state.grid_remove()
        self.file_view.grid()

        for row, fp in enumerate(self.file_queue):
            entry = FileEntry(
                self.files_frame,
                fp,
                on_remove=lambda p=fp: self.remove_file(p),
            )
            entry.grid(row=row, column=0, sticky="ew", pady=(0, 2))
            state, progress, error = self.file_states.get(
                fp, ("Ready", None, None)
            )
            entry.set_state(state, progress, error)
            self.file_widgets.append(entry)
            self.file_widget_map[fp] = entry
        self.update_action_bar()

    def update_action_bar(self):
        count = len(self.file_queue)
        if self.translating:
            return
        self.status_label.configure(
            text="" if count == 0 else f"{count} file{'s' if count != 1 else ''}"
        )
        button_text = "Translate" if count <= 1 else f"Translate {count} files"
        self.translate_btn.configure(text=button_text)
        self.set_translate_available(count > 0)

    def set_file_state(self, filepath, state, progress=None, error=None):
        self.file_states[filepath] = (state, progress, error)
        entry = self.file_widget_map.get(filepath)
        if entry and entry.winfo_exists():
            entry.set_state(state, progress, error)

    def set_all_file_states(self, state):
        for filepath in self.file_queue:
            self.set_file_state(filepath, state)

    def start_translation(self):
        if self.settings_window and self.settings_window.is_open:
            return
        api_key = self.api_key_var.get().strip()
        if not api_key:
            self.show_inline_notice(
                "API key required \u2014 add it in Settings", danger=True
            )
            self.open_settings(focus_api=True)
            return
        if not self.file_queue:
            return

        source_language = self.source_language_var.get()
        target_language = self.target_language_var.get()
        if source_language == target_language:
            messagebox.showinfo(
                "Choose another target",
                "The source and target languages are the same. Choose a different "
                "target language or use Auto-detect.",
                parent=self.master,
            )
            return

        self.save_api_key()

        output_dir = filedialog.askdirectory(
            parent=self.master,
            title="Choose an output folder",
            initialdir=self.config.get("last_output_dir", str(Path.home())),
            mustexist=True,
        )
        if not output_dir:
            return

        self.config["last_output_dir"] = output_dir
        save_config(self.config)

        self.translating = True
        self.cancel_event.clear()
        self.set_translate_available(False)
        self.set_controls_busy(True)
        self.set_all_file_states("Queued")
        self.progress_frame.grid()
        self.progress_bar.set(0)
        self.progress_percent.configure(text="0%")
        self.progress_detail.configure(text="Scanning documents...")
        self.progress_stats.configure(text="")
        self.status_label.configure(text=f"Translating 0 of {len(self.file_queue)}")
        self.cancel_btn.configure(state="normal")
        if source_language == AUTO_DETECT:
            self.detected_source_languages = []
            self.update_swap_state()

        model_id = self.get_selected_model_id()
        files = self.file_queue.copy()

        threading.Thread(
            target=self.run_translation,
            args=(files, output_dir, model_id, api_key, source_language, target_language),
            daemon=True,
        ).start()

    def run_translation(self, files, output_dir, model_id, api_key,
                        source_language, target_language):
        client = OpenAI(api_key=api_key)

        total_paras = 0
        file_para_counts = {}
        scan_errors = {}
        for f in files:
            try:
                ext = os.path.splitext(f)[1].lower()
                if ext == ".pdf":
                    count = scan_pdf_paragraphs(f)
                else:
                    count = scan_pptx_paragraphs(f)
                file_para_counts[f] = count
                total_paras += count
            except Exception as exc:
                file_para_counts[f] = 0
                scan_errors[f] = str(exc)
                self.master.after(
                    0, lambda fp=f, err=str(exc): self.set_file_state(
                        fp, "Failed", error=err
                    )
                )

        if total_paras == 0:
            self.master.after(
                0,
                lambda: (
                    [
                        self.set_file_state(
                            fp,
                            "Failed" if fp in scan_errors else "Ready",
                            error=scan_errors.get(fp),
                        )
                        for fp in files
                    ],
                    self.translation_done(
                        0, 0, True,
                        "No translatable text was found in the selected documents.",
                    ),
                ),
            )
            return

        progress_offset = 0
        translated_paras = 0
        total_input_tokens = 0
        total_output_tokens = 0
        start_time = time.time()
        completed_files = 0
        output_paths = []
        file_results = []
        detected_languages = []

        for file_index, f in enumerate(files, 1):
            if self.cancel_event.is_set():
                break

            filename = os.path.basename(f)
            file_total = file_para_counts.get(f, 0)
            self.master.after(
                0,
                lambda fp=f, idx=file_index, total=len(files): (
                    self.set_file_state(fp, "Translating", 0),
                    self.status_label.configure(text=f"Translating {idx} of {total}"),
                ),
            )
            file_ext = os.path.splitext(filename)[1]
            default_name = os.path.splitext(filename)[0] + "_Translated" + file_ext
            output_path = os.path.join(output_dir, default_name)

            source_for_file = source_language
            if source_language == AUTO_DETECT:
                self.master.after(0, lambda fn=filename: self.progress_detail.configure(
                    text=f"Detecting language: {fn}"
                ))
                detected = None
                try:
                    sample = extract_text_sample(f)
                    detected, detection_in, detection_out = detect_document_language(
                        sample, model_id, client
                    )
                    total_input_tokens += detection_in
                    total_output_tokens += detection_out
                except Exception as exc:
                    print(f"Error reading document for language detection: {exc}")

                if detected:
                    source_for_file = detected
                    detected_languages.append(detected)
                detected_snapshot = detected_languages.copy()
                self.master.after(
                    0,
                    lambda values=detected_snapshot: self.show_detected_languages(values),
                )

            source_status = (
                source_for_file if source_for_file != AUTO_DETECT else "auto-detected language"
            )
            self.master.after(0, lambda fn=filename, lang=source_status: (
                self.progress_detail.configure(text=f"Translating {lang}: {fn}"),
                self.status_label.configure(text=f"Translating to {target_language}"),
            ))

            tokens_before_file_in = total_input_tokens
            tokens_before_file_out = total_output_tokens

            def progress_cb(global_done, in_tok, out_tok, _st=start_time, _tp=total_paras,
                            _tbi=tokens_before_file_in, _tbo=tokens_before_file_out,
                            _offset=progress_offset, _file_total=file_total,
                            _filepath=f):
                nonlocal total_input_tokens, total_output_tokens
                total_input_tokens = _tbi + in_tok
                total_output_tokens = _tbo + out_tok
                elapsed = time.time() - _st
                pct = global_done / _tp if _tp > 0 else 0
                total_tok = total_input_tokens + total_output_tokens
                eta = (elapsed / global_done * (_tp - global_done)) if global_done > 0 else 0
                model_info = MODELS.get(model_id, {})
                cost = (total_input_tokens / 1_000_000) * model_info.get("input_cost", 0) + \
                       (total_output_tokens / 1_000_000) * model_info.get("output_cost", 0)
                pct_int = int(pct * 100)
                file_pct = (
                    min(1.0, max(0.0, (global_done - _offset) / _file_total))
                    if _file_total > 0 else 1.0
                )

                self.master.after(0, lambda p=pct, e=eta, gd=global_done, tp=_tp,
                                  c=cost, it=total_tok, pi=pct_int, fp=_filepath,
                                  fpp=file_pct: (
                    self.progress_bar.set(p),
                    self.progress_percent.configure(text=f"{pi}%"),
                    self.progress_detail.configure(text=f"Paragraph {gd}/{tp}"),
                    self.progress_stats.configure(
                        text=f"About {int(e)}s remaining   \u2022   {it:,} tokens   \u2022   ~${c:.4f}"
                    ),
                    self.set_file_state(fp, "Translating", fpp),
                ))

            try:
                ext = os.path.splitext(f)[1].lower()
                if ext == ".pdf":
                    process_fn = process_pdf
                else:
                    process_fn = process_pptx
                in_tok, out_tok, success, file_paras = process_fn(
                    f, output_path, model_id, client, source_for_file, target_language,
                    progress_callback=progress_cb,
                    cancel_event=self.cancel_event,
                    para_offset=progress_offset,
                )
                total_input_tokens = tokens_before_file_in + in_tok
                total_output_tokens = tokens_before_file_out + out_tok
                if success:
                    translated_paras += file_paras
                    completed_files += 1
                    output_paths.append(output_path)
                    file_results.append({
                        "input_path": str(Path(f).resolve()),
                        "output_path": str(Path(output_path).resolve()),
                    })
                    self.master.after(
                        0, lambda fp=f: self.set_file_state(fp, "Complete", 1.0)
                    )
                else:
                    self.master.after(
                        0, lambda fp=f: self.set_file_state(fp, "Ready")
                    )
            except Exception as e:
                self.master.after(
                    0, lambda fp=f, err=str(e): self.set_file_state(
                        fp, "Failed", error=err
                    )
                )
            finally:
                progress_offset += file_total

        if self.cancel_event.is_set():
            for filepath in files:
                state = self.file_states.get(filepath, ("Ready", None, None))[0]
                if state == "Queued":
                    self.master.after(
                        0, lambda fp=filepath: self.set_file_state(fp, "Ready")
                    )
            self.master.after(0, lambda: self.translation_done(
                total_input_tokens, total_output_tokens, False, "Translation cancelled.",
                None, None
            ))
        else:
            model_info = MODELS.get(model_id, {})
            cost = (total_input_tokens / 1_000_000) * model_info.get("input_cost", 0) + \
                   (total_output_tokens / 1_000_000) * model_info.get("output_cost", 0)
            elapsed = time.time() - start_time
            result_info = {
                "completed_files": completed_files,
                "total_files": len(files),
                "paragraphs": translated_paras,
                "tokens": total_input_tokens + total_output_tokens,
                "cost": cost,
                "elapsed": int(elapsed),
                "output_dir": output_dir,
                "input_paths": [str(Path(path).resolve()) for path in files],
                "output_paths": output_paths,
                "file_results": file_results,
                "model_id": model_id,
                "source_languages": (
                    detected_languages if source_language == AUTO_DETECT else [source_language]
                ),
                "target_language": target_language,
            }
            # Show 100% briefly before showing completion window
            self.master.after(0, lambda: (
                self.progress_bar.set(1.0),
                self.progress_percent.configure(text="100%"),
                self.progress_detail.configure(text="Translation complete"),
                self.status_label.configure(text=f"Translated {completed_files} of {len(files)}"),
            ))
            self.master.after(0, lambda ri=result_info: self.translation_done(
                total_input_tokens, total_output_tokens, True, "Translation complete!",
                ri.get("output_dir"), ri.get("output_paths"), ri
            ))

    def translation_done(self, input_tokens, output_tokens, success, message,
                         output_dir=None, output_paths=None, result_info=None):
        self.translating = False
        self.set_controls_busy(False)
        self.update_action_bar()
        self.progress_frame.grid_remove()

        if success and result_info:
            self.status_label.configure(
                text=f"{result_info['completed_files']} of "
                     f"{result_info['total_files']} files complete"
            )
            self.add_history_record(result_info)
            self.show_completion_panel(result_info)
        elif not success:
            self.show_inline_notice(message, danger=True)
            self.status_label.configure(text="Translation cancelled")
        else:
            messagebox.showinfo("Information", message, parent=self.master)
            self.update_action_bar()

    def show_completion_panel(self, info):
        if self.completion_panel and self.completion_panel.overlay.winfo_exists():
            self.completion_panel.overlay.destroy()
        self.completion_panel = CompletionPanel(self, info)

    def open_in_folder(self, folder_path):
        try:
            open_path(folder_path)
        except OSError as exc:
            messagebox.showerror(
                "Unable to open folder", str(exc), parent=self.master
            )

    def open_file(self, file_path):
        try:
            open_path(file_path)
        except OSError as exc:
            messagebox.showerror(
                "Unable to open document", str(exc), parent=self.master
            )

    def reveal_file(self, file_path):
        try:
            reveal_path(file_path)
        except OSError as exc:
            messagebox.showerror(
                "Unable to reveal document", str(exc), parent=self.master
            )

    def cancel_translation(self):
        self.cancel_event.set()
        self.cancel_btn.configure(state="disabled")
        self.progress_detail.configure(text="Cancelling...")


if __name__ == "__main__":
    enable_windows_dpi_awareness()
    app = DragDropCTk()
    PPTTranslatorApp(app)
    app.mainloop()
