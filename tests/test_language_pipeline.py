import tempfile
import unittest
from pathlib import Path
from types import SimpleNamespace

from pptx import Presentation
from pptx.util import Inches
import pymupdf as fitz

import pptxtranslator as app


class FakeCompletions:
    def __init__(self):
        self.system_prompts = []

    def create(self, model, messages, temperature=0, **kwargs):
        system_prompt = messages[0]["content"]
        self.system_prompts.append(system_prompt)
        if system_prompt.startswith("Identify the dominant language"):
            content = "Japanese"
        else:
            segments = messages[-1]["content"].split(app.SEP_TOKEN)
            content = app.SEP_TOKEN.join(f"English: {segment}" for segment in segments)
        return SimpleNamespace(
            choices=[SimpleNamespace(message=SimpleNamespace(content=content))],
            usage=SimpleNamespace(prompt_tokens=12, completion_tokens=6),
        )


class FakeClient:
    def __init__(self):
        self.chat = SimpleNamespace(completions=FakeCompletions())


class LanguagePipelineTests(unittest.TestCase):
    def test_detection_response_is_canonicalized(self):
        self.assertEqual(
            app.canonicalize_detected_language("Traditional Chinese"),
            "Chinese (Traditional)",
        )
        self.assertEqual(
            app.canonicalize_detected_language("Brazilian Portuguese"),
            "Portuguese (Brazil)",
        )

    def test_pptx_text_boxes_and_tables_are_translated(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            source_path = Path(temp_dir) / "source.pptx"
            output_path = Path(temp_dir) / "translated.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            textbox.text_frame.text = "日本語の見出し"
            table = slide.shapes.add_table(1, 1, Inches(1), Inches(2), Inches(4), Inches(1)).table
            table.cell(0, 0).text = "こんにちは"
            presentation.save(source_path)

            self.assertEqual(app.scan_pptx_paragraphs(source_path), 2)
            self.assertIn("日本語", app.extract_text_sample(source_path))

            client = FakeClient()
            detected, _, _ = app.detect_document_language(
                app.extract_text_sample(source_path), "test-model", client
            )
            self.assertEqual(detected, "Japanese")

            _, _, success, count = app.process_pptx(
                source_path,
                output_path,
                "test-model",
                client,
                "Japanese",
                "English",
            )
            self.assertTrue(success)
            self.assertEqual(count, 2)

            translated = Presentation(output_path)
            paragraphs = [
                "".join(run.text for run in paragraph.runs)
                for paragraph in app.iter_pptx_paragraphs(translated)
            ]
            self.assertIn("English: 日本語の見出し", paragraphs)
            self.assertIn("English: こんにちは", paragraphs)
            self.assertTrue(
                any("into English" in prompt for prompt in client.chat.completions.system_prompts)
            )

    def test_pdf_complex_script_uses_unicode_text_insertion(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            source_path = Path(temp_dir) / "source.pdf"
            output_path = Path(temp_dir) / "translated.pdf"

            document = fitz.open()
            page = document.new_page()
            page.insert_text((72, 100), "Hello world", fontsize=12)
            document.save(source_path)
            document.close()

            completions = SimpleNamespace(
                create=lambda **_kwargs: SimpleNamespace(
                    choices=[
                        SimpleNamespace(
                            message=SimpleNamespace(content="வணக்கம் உலகம்")
                        )
                    ],
                    usage=SimpleNamespace(prompt_tokens=3, completion_tokens=3),
                )
            )
            client = SimpleNamespace(chat=SimpleNamespace(completions=completions))

            _, _, success, count = app.process_pdf(
                source_path,
                output_path,
                "test-model",
                client,
                "English",
                "Tamil",
            )
            self.assertTrue(success)
            self.assertEqual(count, 1)

            translated = fitz.open(output_path)
            try:
                output_text = "".join(page.get_text() for page in translated)
            finally:
                translated.close()
            self.assertIn("வணக்கம்", output_text)


if __name__ == "__main__":
    unittest.main()
